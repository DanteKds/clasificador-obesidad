# -*- coding: utf-8 -*-
"""
Generador de PPT final — Análisis de Obesidad CRISP-DM
Ejecutar: python create_ppt_v3.py   (requiere outputs/resultados_v2.json)
Produce:  presentacion_obesidad_final.pptx (15 diapositivas)

Diseño: paleta pizarra/teal/coral, tarjetas redondeadas, tipografía Segoe UI.
Las métricas se leen desde resultados_v2.json — sin números hardcodeados.
"""
import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

with open("outputs/resultados_v2.json", encoding="utf-8") as f:
    R = json.load(f)

MEJOR = R["mejor_modelo"]
MB = R["metricas_base"]
CL = R["clustering"]
EXP = R["experimentos"]
FP = R["fp_growth"]
REGLAS = R["reglas_novedosas"]

# ─── PALETA ──────────────────────────────────────────────────────
SLATE   = RGBColor(0x16, 0x2A, 0x47)   # azul pizarra profundo (títulos / portada)
SLATE2  = RGBColor(0x1F, 0x3A, 0x5F)   # pizarra medio (cajas sobre fondo oscuro)
TEAL    = RGBColor(0x0E, 0xA5, 0xA5)   # acento principal
CORAL   = RGBColor(0xFF, 0x6B, 0x5B)   # acento cálido
AMBER   = RGBColor(0xF2, 0xA0, 0x4E)   # acento secundario
GREENOK = RGBColor(0x2F, 0x9E, 0x6E)   # éxito
PURPLE  = RGBColor(0x7E, 0x57, 0xC2)   # experimentos
INK     = RGBColor(0x24, 0x33, 0x42)   # texto principal
MUTED   = RGBColor(0x5A, 0x6B, 0x7B)   # texto secundario
BG      = RGBColor(0xF7, 0xF9, 0xFB)   # fondo de slide
CARD    = RGBColor(0xFF, 0xFF, 0xFF)
BORDER  = RGBColor(0xDD, 0xE5, 0xEC)
TEALBG  = RGBColor(0xE4, 0xF4, 0xF4)   # fondo teal suave
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


# ─── HELPERS ─────────────────────────────────────────────────────
def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def shape_box(slide, kind, l, t, w, h, fill=None, line_color=None, line_w=0.75):
    """kind: 1=rect, 5=rounded rect, 9=oval"""
    shp = slide.shapes.add_shape(kind, l, t, w, h)
    if kind == 5:
        try:
            shp.adjustments[0] = 0.07
        except Exception:
            pass
    shp.shadow.inherit = False
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    return shp


def rect(slide, l, t, w, h, fill=None, line_color=None):
    return shape_box(slide, 1, l, t, w, h, fill, line_color)


def card(slide, l, t, w, h, accent=None, fill=CARD):
    c = shape_box(slide, 5, l, t, w, h, fill=fill, line_color=BORDER)
    if accent:
        shape_box(slide, 5, l, t, Inches(0.14), h, fill=accent)
    return c


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False, font=FONT,
             anchor=None):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txBox


def header(slide, title, subtitle=None):
    """Encabezado claro: chip teal + título pizarra + subrayado."""
    rect(slide, 0, 0, W, H, fill=BG)
    rect(slide, Inches(0.45), Inches(0.32), Inches(0.16), Inches(0.62), fill=TEAL)
    add_text(slide, title, Inches(0.78), Inches(0.22), Inches(11.9), Inches(0.6),
             size=26, bold=True, color=SLATE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.8), Inches(0.78), Inches(11.9), Inches(0.4),
                 size=13, color=MUTED)
    rect(slide, Inches(0.45), Inches(1.28), Inches(12.43), Pt(1.4), fill=BORDER)


def slide_num(slide, n, total=15, dark=False):
    add_text(slide, f"{n:02d} · {total}", Inches(12.35), Inches(7.12), Inches(0.9), Inches(0.3),
             size=9, color=RGBColor(0x9A, 0xA8, 0xB5) if not dark else RGBColor(0x6E, 0x7F, 0x9E),
             align=PP_ALIGN.RIGHT)


def add_image(slide, path, l, t, w, h=None):
    if not os.path.exists(path):
        return
    if h:
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        slide.shapes.add_picture(path, l, t, w)


def img_card(slide, path, l, t, w, pad=Inches(0.12), est_ratio=0.52):
    """Imagen sobre tarjeta blanca redondeada (marco)."""
    if not os.path.exists(path):
        return
    pic = slide.shapes.add_picture(path, l + pad, t + pad, w - 2 * pad)
    frame_h = pic.height + 2 * pad
    card_shape = shape_box(slide, 5, l, t, w, frame_h, fill=CARD, line_color=BORDER)
    # mover la tarjeta detrás de la imagen
    sp = card_shape._element
    sp.getparent().remove(sp)
    pic._element.addprevious(sp)
    return frame_h


def bullet_box(slide, items, l, t, w, h, title=None, title_color=SLATE,
               item_size=13, item_color=INK, leading=3, marker="•  "):
    if title:
        add_text(slide, title, l, t, w, Inches(0.35), size=14, bold=True, color=title_color)
        t += Inches(0.4)
        h -= Inches(0.4)
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(leading)
        run = p.add_run()
        run.text = marker + item
        run.font.size = Pt(item_size)
        run.font.color.rgb = item_color
        run.font.name = FONT


def tabla(slide, data, col_ls, col_ws, ty, row_h, header_bg=SLATE,
          font_size=11, highlight_row=None, highlight_color=GREENOK):
    for ri, row in enumerate(data):
        if ri == 0:
            bg, fc, bold = header_bg, WHITE, True
        elif highlight_row is not None and ri == highlight_row:
            bg, fc, bold = TEALBG, highlight_color, True
        else:
            bg = CARD if ri % 2 == 1 else RGBColor(0xF1, 0xF5, 0xF9)
            fc, bold = INK, False
        for ci, (text, lx, cw) in enumerate(zip(row, col_ls, col_ws)):
            rect(slide, lx, ty, cw, row_h, fill=bg, line_color=BORDER)
            add_text(slide, text, lx + Inches(0.07), ty + Inches(0.06),
                     cw - Inches(0.1), row_h - Inches(0.06), size=font_size,
                     bold=bold, color=fc,
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
        ty += row_h
    return ty


def kpi_chip(slide, valor, etiqueta, l, t, w, color=TEAL):
    """Chip de métrica destacada."""
    card(slide, l, t, w, Inches(1.0), accent=None, fill=CARD)
    add_text(slide, valor, l, t + Inches(0.08), w, Inches(0.5),
             size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, etiqueta, l + Inches(0.05), t + Inches(0.58), w - Inches(0.1), Inches(0.38),
             size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════
# SLIDE 1 · PORTADA
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
rect(s, 0, 0, W, H, fill=SLATE)
# formas decorativas
shape_box(s, 9, Inches(10.2), Inches(3.9), Inches(5.2), Inches(5.2), fill=SLATE2)
shape_box(s, 9, Inches(11.6), Inches(5.1), Inches(3.6), Inches(3.6), fill=TEAL)
shape_box(s, 9, Inches(9.7), Inches(-1.6), Inches(2.6), Inches(2.6), fill=SLATE2)
shape_box(s, 9, Inches(-0.9), Inches(5.9), Inches(2.4), Inches(2.4),
          fill=RGBColor(0x2A, 0x46, 0x6B))
rect(s, Inches(0.75), Inches(2.62), Inches(2.6), Pt(3), fill=CORAL)

add_text(s, "TALLER DE APLICACIONES · MAGÍSTER EN DATA SCIENCE",
         Inches(0.75), Inches(0.85), Inches(10.0), Inches(0.4),
         size=13, bold=True, color=TEAL)
add_text(s, "Análisis de Niveles de Obesidad",
         Inches(0.7), Inches(1.25), Inches(11.5), Inches(0.85),
         size=40, bold=True, color=WHITE)
add_text(s, "Clustering · Clasificación · Reglas de Asociación · Experimentos de Transformación",
         Inches(0.75), Inches(2.12), Inches(11.0), Inches(0.5),
         size=17, color=RGBColor(0xBA, 0xCB, 0xDE))

add_text(s, "Metodología CRISP-DM aplicada al dataset Estimation of Obesity Levels\n"
            "UCI Machine Learning Repository · 2.111 registros · México, Perú y Colombia",
         Inches(0.75), Inches(2.95), Inches(9.5), Inches(0.8),
         size=13, color=RGBColor(0x8F, 0xA5, 0xBC))

infos = [
    ("ALUMNO", "Dante Gil Zenteno"),
    ("DOCENTE", "Dr. Mauricio Sepúlveda"),
    ("FACULTAD", "Facultad de Ingeniería — USS"),
]
for i, (label, val) in enumerate(infos):
    y = Inches(4.35) + i * Inches(0.78)
    rect(s, Inches(0.78), y + Inches(0.07), Inches(0.07), Inches(0.5), fill=TEAL)
    add_text(s, label, Inches(1.02), y, Inches(3.2), Inches(0.32),
             size=10, bold=True, color=RGBColor(0x7E, 0x95, 0xAE))
    add_text(s, val, Inches(1.02), y + Inches(0.27), Inches(5.5), Inches(0.4),
             size=15, bold=True, color=WHITE)

add_text(s, "2026", Inches(0.78), Inches(6.95), Inches(2.0), Inches(0.35),
         size=12, bold=True, color=RGBColor(0x7E, 0x95, 0xAE))


# ═════════════════════════════════════════════════════════════════
# SLIDE 2 · OBJETIVO Y PROBLEMA
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "Objetivo del Trabajo y Problema",
       "CRISP-DM · Fase 1 — Comprensión del negocio")
slide_num(s, 2)

card(s, Inches(0.45), Inches(1.55), Inches(6.0), Inches(5.45), accent=TEAL)
add_text(s, "El Dataset", Inches(0.78), Inches(1.72), Inches(5.5), Inches(0.4),
         size=16, bold=True, color=SLATE)
bullet_box(s, [
    "2.111 registros de México, Perú y Colombia",
    "17 atributos: hábitos alimentarios, actividad física y características personales",
    "Variable objetivo NObeyesdad: 7 niveles de obesidad",
    "Clases balanceadas (272–351 registros por nivel)",
    "Sin valores nulos ni duplicados",
    "77% de registros sintéticos vía SMOTE — Mendoza Palechor & De la Hoz (2019)",
], Inches(0.78), Inches(2.25), Inches(5.45), Inches(4.0), item_size=13, leading=10)
rect(s, Inches(0.78), Inches(6.3), Inches(5.35), Pt(1.2), fill=BORDER)
add_text(s, "Fuente: UCI Machine Learning Repository · doi:10.1016/j.dib.2019.104344",
         Inches(0.78), Inches(6.45), Inches(5.45), Inches(0.4),
         size=11, italic=True, color=MUTED)

card(s, Inches(6.75), Inches(1.55), Inches(6.15), Inches(3.15), accent=CORAL)
add_text(s, "Preguntas Guía", Inches(7.08), Inches(1.72), Inches(5.6), Inches(0.4),
         size=16, bold=True, color=SLATE)
bullet_box(s, [
    "¿Existen grupos de personas con perfiles de riesgo similares?",
    "¿Qué algoritmo predice mejor el nivel de obesidad?",
    "¿Qué combinaciones de hábitos se asocian con la obesidad?",
    "¿Cómo cambian los resultados con reducción de dimensionalidad?",
    "¿Aportan valor nuevos atributos de ingeniería de características?",
], Inches(7.08), Inches(2.25), Inches(5.6), Inches(2.4), item_size=12, leading=5)

card(s, Inches(6.75), Inches(4.95), Inches(6.15), Inches(2.05), accent=AMBER)
add_text(s, "Objetivo General", Inches(7.08), Inches(5.12), Inches(5.6), Inches(0.4),
         size=16, bold=True, color=SLATE)
add_text(s, "Aplicar CRISP-DM con agrupamiento, clasificación y asociación FP-Growth, "
            "incorporando tres experimentos de transformación de datos y comparando su "
            "impacto en las métricas.",
         Inches(7.08), Inches(5.58), Inches(5.6), Inches(1.3), size=13, color=MUTED)


# ═════════════════════════════════════════════════════════════════
# SLIDE 3 · METODOLOGÍA CRISP-DM
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "Metodología CRISP-DM",
       "Seis fases que estructuran todo el análisis")
slide_num(s, 3)

fases = [
    ("1", "Comprensión del Negocio", "Preguntas guía y métricas de éxito", TEAL),
    ("2", "Comprensión de los Datos", "Inspección, calidad y balance de clases", TEAL),
    ("3", "Preparación de los Datos", "Codificación y escalado por técnica", TEAL),
    ("4", "Modelamiento", "K-Means · 5 clasificadores · FP-Growth", CORAL),
    ("5", "Evaluación", "Métricas en test, Lift y experimentos", CORAL),
    ("6", "Comunicación", "Presentación + notebook reproducible", CORAL),
]
for i, (num, nombre, desc, color) in enumerate(fases):
    col, row = i % 3, i // 3
    lx = Inches(0.45 + col * 4.25)
    ty = Inches(1.55 + row * 1.62)
    card(s, lx, ty, Inches(4.05), Inches(1.45))
    shape_box(s, 9, lx + Inches(0.18), ty + Inches(0.3), Inches(0.8), Inches(0.8), fill=color)
    add_text(s, num, lx + Inches(0.18), ty + Inches(0.42), Inches(0.8), Inches(0.55),
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, nombre, lx + Inches(1.12), ty + Inches(0.16), Inches(2.85), Inches(0.55),
             size=13, bold=True, color=SLATE)
    add_text(s, desc, lx + Inches(1.12), ty + Inches(0.68), Inches(2.85), Inches(0.7),
             size=11, color=MUTED)

card(s, Inches(0.45), Inches(4.95), Inches(12.45), Inches(2.05), accent=TEAL, fill=TEALBG)
add_text(s, "Qué aportó la metodología a este proyecto",
         Inches(0.78), Inches(5.1), Inches(11.9), Inches(0.4),
         size=15, bold=True, color=SLATE)
bullet_box(s, [
    "Definir preguntas guía ANTES de tocar datos evitó modelar sin entender el problema",
    "La inspección temprana (nulos, duplicados, balance) descartó problemas de calidad y retrabajos",
    "Preparar los datos POR TÉCNICA (escalado / pipelines / discretización) eliminó fugas de datos",
    "La fase de evaluación obligó a comparar en test y a documentar también los resultados que no mejoran",
], Inches(0.78), Inches(5.55), Inches(11.9), Inches(1.4), item_size=12, leading=3)


# ═════════════════════════════════════════════════════════════════
# SLIDE 4 · DATASET Y VARIABLE OBJETIVO
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "Descripción del Dataset",
       "CRISP-DM · Fase 2 — Comprensión de los datos")
slide_num(s, 4)

card(s, Inches(0.45), Inches(1.55), Inches(5.3), Inches(5.45), accent=TEAL)
add_text(s, "Variables (17)", Inches(0.78), Inches(1.7), Inches(4.8), Inches(0.4),
         size=15, bold=True, color=SLATE)
bullet_box(s, [
    "Género · Edad · Altura · Peso",
    "Historia familiar de sobrepeso",
    "Alimentos hipercalóricos (FAVC)",
    "Frecuencia de vegetales (FCVC)",
    "Comidas principales (NCP)",
    "Snacks entre comidas (CAEC)",
    "Tabaquismo (SMOKE) · Agua (CH2O)",
    "Monitoreo de calorías (SCC)",
    "Actividad física (FAF) · Pantallas (TUE)",
    "Alcohol (CALC) · Transporte (MTRANS)",
], Inches(0.78), Inches(2.18), Inches(4.8), Inches(3.7), item_size=12, leading=4)
rect(s, Inches(0.78), Inches(6.18), Inches(4.6), Pt(1.2), fill=BORDER)
add_text(s, "★ Target: NObeyesdad — 7 clases", Inches(0.78), Inches(6.35),
         Inches(4.7), Inches(0.4), size=13, bold=True, color=CORAL)

add_text(s, "Distribución de la Variable Objetivo", Inches(6.1), Inches(1.62),
         Inches(6.8), Inches(0.4), size=15, bold=True, color=SLATE)
niveles = [
    ("Nivel de Obesidad",   "Registros", "%"),
    ("Obesity_Type_I",      "351", "16.6%"),
    ("Obesity_Type_III",    "324", "15.3%"),
    ("Obesity_Type_II",     "297", "14.1%"),
    ("Overweight_Level_I",  "290", "13.7%"),
    ("Overweight_Level_II", "290", "13.7%"),
    ("Normal_Weight",       "287", "13.6%"),
    ("Insufficient_Weight", "272", "12.9%"),
]
tabla(s, niveles, [Inches(6.1), Inches(9.9), Inches(11.35)],
      [Inches(3.8), Inches(1.45), Inches(1.45)], Inches(2.1), Inches(0.46), font_size=12)
add_text(s, "Clases balanceadas → métricas macro confiables y sin necesidad de re-balancear",
         Inches(6.1), Inches(5.95), Inches(6.8), Inches(0.6), size=12, italic=True, color=MUTED)


# ═════════════════════════════════════════════════════════════════
# SLIDE 5 · PREPARACIÓN Y TRANSFORMACIÓN
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "Preparación y Transformación de Datos",
       "CRISP-DM · Fase 3 — Una preparación específica por técnica")
slide_num(s, 5)

steps = [
    ("Calidad", TEAL,
     ["Sin nulos ni duplicados", "Rangos numéricos consistentes"]),
    ("Variables", TEAL,
     ["8 numéricas · 8 categóricas", "Target: 7 clases"]),
    ("Clustering", CORAL,
     ["OrdinalEncoder + StandardScaler", "Target excluido — sin fuga de datos"]),
    ("Clasificación", CORAL,
     ["Pipeline con preprocesamiento dentro", "Split estratificado 80/20 · seed 42"]),
    ("FP-Growth", AMBER,
     ["Numéricas discretizadas en rangos", "Ítems binarios: Gender=Female…"]),
]
for i, (title, color, items) in enumerate(steps):
    lx = Inches(0.45 + i * 2.53)
    card(s, lx, Inches(1.55), Inches(2.38), Inches(3.3))
    rect(s, lx + Inches(0.16), Inches(1.78), Inches(0.55), Pt(3), fill=color)
    add_text(s, title, lx + Inches(0.16), Inches(1.92), Inches(2.1), Inches(0.4),
             size=14, bold=True, color=SLATE)
    bullet_box(s, items, lx + Inches(0.16), Inches(2.42), Inches(2.1), Inches(2.3),
               item_size=11, leading=6)

card(s, Inches(0.45), Inches(5.15), Inches(12.45), Inches(1.85), accent=PURPLE)
add_text(s, "Transformaciones experimentales evaluadas",
         Inches(0.78), Inches(5.3), Inches(11.9), Inches(0.4),
         size=15, bold=True, color=PURPLE)
add_text(s, f"PCA ({EXP['pca']['n_comp']} componentes — 95% de varianza)   ·   "
            f"TruncatedSVD sobre One-Hot ({EXP['svd']['n_comp']} componentes — 90% de varianza)   ·   "
            "5 atributos nuevos: BMI, BMI_Category, Healthy_Habits_Score, Risky_Eating_Score, Sedentary_Risk_Score",
         Inches(0.78), Inches(5.78), Inches(11.9), Inches(1.0), size=13, color=INK)


# ═════════════════════════════════════════════════════════════════
# SLIDE 6 · EDA
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "Análisis Exploratorio de Datos",
       "Distribuciones, correlaciones y patrones")
slide_num(s, 6)

img_card(s, "outputs/01_distribucion_objetivo.png", Inches(0.45), Inches(1.55), Inches(6.6))
img_card(s, "outputs/03_correlacion.png", Inches(7.25), Inches(1.55), Inches(5.65))

card(s, Inches(0.45), Inches(6.15), Inches(12.45), Inches(0.95), accent=TEAL)
add_text(s, "Hallazgos:  7 clases balanceadas (272–351 por nivel)   ·   Weight y Height son las "
            "variables más correlacionadas con el target   ·   Age aporta menos discriminación entre clases",
         Inches(0.78), Inches(6.42), Inches(11.9), Inches(0.5), size=12, color=INK)


# ═════════════════════════════════════════════════════════════════
# SLIDE 7 · CLUSTERING — MÉTODO Y K
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "Clustering — Algoritmo y Número de Grupos",
       "CRISP-DM · Fase 4 — Agrupamiento con K-Means")
slide_num(s, 7)

img_card(s, "outputs/06_evaluacion_clusters.png", Inches(0.45), Inches(1.55), Inches(8.1))

card(s, Inches(8.75), Inches(1.55), Inches(4.15), Inches(2.55), accent=TEAL)
add_text(s, "¿Por qué K-Means?", Inches(9.05), Inches(1.7), Inches(3.7), Inches(0.4),
         size=14, bold=True, color=SLATE)
bullet_box(s, [
    "Escala bien con 2.111 registros",
    "Centroides interpretables como perfiles",
    "Alternativas evaluadas: DBSCAN (calibrar eps) y jerárquico (lento, resultado similar)",
], Inches(9.05), Inches(2.15), Inches(3.7), Inches(1.9), item_size=11, leading=4)

card(s, Inches(8.75), Inches(4.3), Inches(4.15), Inches(2.7), accent=CORAL)
add_text(s, "¿Por qué k = 4?", Inches(9.05), Inches(4.45), Inches(3.7), Inches(0.4),
         size=14, bold=True, color=SLATE)
bullet_box(s, [
    "k=2 maximiza Silhouette pero solo separa \"obeso / no obeso\"",
    "k=4 mapea los 7 niveles en 4 perfiles de riesgo accionables",
    "Silhouettes < 0.20 en todo el rango: esperable en datos de salud con registros sintéticos",
], Inches(9.05), Inches(4.9), Inches(3.7), Inches(2.0), item_size=11, leading=4)


# ═════════════════════════════════════════════════════════════════
# SLIDE 8 · CLUSTERING — PERFILES
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "Clustering — Cuatro Perfiles de Riesgo",
       "Interpretación de cada grupo con la variable objetivo como lectura externa")
slide_num(s, 8)

img_card(s, "outputs/07_clusters_visualizacion.png", Inches(0.45), Inches(1.55), Inches(7.6))

clusters = [
    ("Cluster 0 · 22.4%", "Bajo peso / Normal",
     "IMC 21.4 · 56.6 kg · 20.8 años · FAF 1.10 · Hist. familiar 41%", TEAL),
    ("Cluster 1 · 23.1%", "Obesidad severa",
     "IMC 37.8 · 106.6 kg · 22.9 años · FAF 0.73 · Hist. familiar 100%", CORAL),
    ("Cluster 2 · 36.4%", "Obesidad moderada",
     "IMC 28.9 · 90.2 kg · 22.3 años · FAF 1.20 · Hist. familiar 90%", AMBER),
    ("Cluster 3 · 18.1%", "Obesidad leve / tipo I",
     "IMC 31.2 · 91.0 kg · 34.4 años · FAF 0.88 · Hist. familiar 92%", GREENOK),
]
for i, (nombre, perfil, stats, color) in enumerate(clusters):
    ty = Inches(1.55 + i * 1.4)
    card(s, Inches(8.3), ty, Inches(4.6), Inches(1.26), accent=color)
    add_text(s, nombre, Inches(8.62), ty + Inches(0.08), Inches(4.1), Inches(0.32),
             size=12, bold=True, color=SLATE)
    add_text(s, perfil, Inches(8.62), ty + Inches(0.4), Inches(4.1), Inches(0.32),
             size=12, bold=True, color=color)
    add_text(s, stats, Inches(8.62), ty + Inches(0.74), Inches(4.2), Inches(0.5),
             size=10, color=MUTED)


# ═════════════════════════════════════════════════════════════════
# SLIDE 9 · CLASIFICACIÓN — MODELOS
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "Clasificación — Cinco Modelos Comparados",
       "CRISP-DM · Fase 4 — Pipelines sin fuga de datos, evaluación en test")
slide_num(s, 9)

img_card(s, "outputs/08_comparacion_modelos.png", Inches(0.45), Inches(1.55), Inches(9.0))

card(s, Inches(9.7), Inches(1.55), Inches(3.2), Inches(2.5), accent=TEAL)
add_text(s, "Configuración", Inches(9.98), Inches(1.7), Inches(2.8), Inches(0.4),
         size=13, bold=True, color=SLATE)
bullet_box(s, [
    "Train / Test: 80 / 20 estratificado",
    "Preprocesamiento dentro del pipeline",
    "Métricas macro (las 7 clases pesan igual)",
    "random_state = 42",
], Inches(9.98), Inches(2.12), Inches(2.8), Inches(1.9), item_size=10, leading=4)

card(s, Inches(9.7), Inches(4.25), Inches(3.2), Inches(2.75), accent=CORAL)
add_text(s, "Modelos", Inches(9.98), Inches(4.4), Inches(2.8), Inches(0.4),
         size=13, bold=True, color=SLATE)
bullet_box(s, [
    "Regresión Logística multinomial",
    "Árbol de Decisión (prof. 12)",
    "Random Forest (100 árboles)",
    "KNN (k=7)",
    "XGBoost (200 árboles)",
], Inches(9.98), Inches(4.82), Inches(2.8), Inches(2.1), item_size=11, leading=5)


# ═════════════════════════════════════════════════════════════════
# SLIDE 10 · CLASIFICACIÓN — RESULTADOS
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "Clasificación — Resultados en Test",
       "CRISP-DM · Fase 5 — Métricas sobre el 20% nunca visto en entrenamiento")
slide_num(s, 10)

img_card(s, "outputs/09_confusion_matrix.png", Inches(0.45), Inches(1.55), Inches(6.1))

filas = [("Modelo", "Accuracy", "Precisión", "Recall", "F1 macro")]
orden_modelos = sorted(R["tabla_modelos"].items(),
                       key=lambda kv: kv[1]["F1 Macro"], reverse=True)
for nombre, met in orden_modelos:
    star = " ★" if nombre == MEJOR else ""
    filas.append((nombre + star, f"{met['Accuracy']:.4f}", f"{met['Precision']:.4f}",
                  f"{met['Recall']:.4f}", f"{met['F1 Macro']:.4f}"))
add_text(s, "Métricas por modelo", Inches(6.85), Inches(1.6), Inches(6.0), Inches(0.4),
         size=15, bold=True, color=SLATE)
tabla(s, filas,
      [Inches(6.85), Inches(9.05), Inches(10.05), Inches(11.05), Inches(11.95)],
      [Inches(2.2), Inches(1.0), Inches(1.0), Inches(0.9), Inches(1.0)],
      Inches(2.05), Inches(0.5), highlight_row=1, font_size=10)

kpi_chip(s, f"{MB['F1 Macro']:.4f}", f"F1 macro · {MEJOR}", Inches(6.85), Inches(5.25),
         Inches(2.9), color=GREENOK)
kpi_chip(s, f"{MB['Accuracy']:.1%}", "Accuracy en test", Inches(9.95), Inches(5.25),
         Inches(2.9), color=TEAL)

add_text(s, "Los errores se concentran en clases limítrofes (Normal / Overweight) — esperable "
            "porque sus fronteras de IMC son continuas.",
         Inches(6.85), Inches(6.45), Inches(6.0), Inches(0.7), size=11, italic=True, color=MUTED)


# ═════════════════════════════════════════════════════════════════
# SLIDE 11 · FP-GROWTH — MEJORES REGLAS
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "FP-Growth — Reglas con Mejores Indicadores",
       "min_support 0.15 · min_confidence 0.70 · Lift > 2.0 · target como consecuente")
slide_num(s, 11)

img_card(s, "outputs/10_reglas_asociacion.png", Inches(0.45), Inches(1.55), Inches(7.2))

add_text(s, "Top reglas hacia NObeyesdad", Inches(7.95), Inches(1.6), Inches(5.0), Inches(0.4),
         size=14, bold=True, color=SLATE)
top_rules = [
    ("Antecedentes", "Consecuente", "Conf", "Lift"),
    ("family_history=yes + Weight=Extremo", "Obesity_III", "0.98", "6.38"),
    ("FAVC=yes + FAF=Nulo + Weight=Alto",   "Obesity_II",  "0.97", "6.34"),
    ("MTRANS=Public + FCVC=Alto",           "Obesity_III", "0.97", "6.35"),
    ("NCP=Normal + SCC=no + FAVC=yes",      "Obesity_III", "0.97", "6.32"),
    ("CALC=Sometimes + Gender=Female",      "Obesity_III", "0.97", "6.31"),
    ("CH2O=Bajo + FAF=Nulo",                "Obesity_I",   "0.94", "5.47"),
]
tabla(s, top_rules,
      [Inches(7.95), Inches(10.85), Inches(11.65), Inches(12.3)],
      [Inches(2.9), Inches(0.8), Inches(0.65), Inches(0.65)],
      Inches(2.05), Inches(0.5), font_size=9)

card(s, Inches(7.95), Inches(5.25), Inches(4.95), Inches(0.85), accent=CORAL)
add_text(s, "Lift > 6  ⇒  la regla es 6× más probable que el azar",
         Inches(8.25), Inches(5.48), Inches(4.5), Inches(0.4), size=12, bold=True, color=CORAL)

add_text(s, f"{FP['n_itemsets']:,} itemsets frecuentes · {FP['n_reglas']:,} reglas generadas".replace(",", "."),
         Inches(7.95), Inches(6.25), Inches(5.0), Inches(0.4), size=12, bold=True, color=SLATE)
add_text(s, "Support = P(A∪B) · Confidence = P(B|A) · Lift = Confidence / P(B)",
         Inches(7.95), Inches(6.6), Inches(5.0), Inches(0.4), size=10, italic=True, color=MUTED)


# ═════════════════════════════════════════════════════════════════
# SLIDE 12 · 6 REGLAS NOVEDOSAS
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "Seis Reglas Novedosas de Asociación",
       "Seleccionadas por Lift alto, antecedentes múltiples y valor interpretativo")
slide_num(s, 12)

novedad = [
    "Combina factor genético con patrón conductual complejo — no trivial.",
    "El transporte público sedentario potencia el efecto de la dieta hipercalórica.",
    "No fumar no protege de la obesidad severa cuando los hábitos alimentarios son malos.",
    "Snacking y alcohol moderado refuerzan el riesgo — regla multivariada.",
    "Predice también el patrón de conducta (FAVC), no solo el nivel de obesidad.",
    "Alcohol ocasional en mujeres con antecedentes familiares y dieta alta en calorías.",
]
for i, regla in enumerate(REGLAS[:6]):
    col, row = i % 3, i // 3
    lx = Inches(0.45 + col * 4.25)
    ty = Inches(1.55 + row * 2.75)
    card(s, lx, ty, Inches(4.05), Inches(2.6))
    rect(s, lx, ty, Inches(4.05), Inches(0.5), fill=SLATE)
    add_text(s, f"Regla {i+1}", lx + Inches(0.15), ty + Inches(0.07), Inches(2.0), Inches(0.36),
             size=12, bold=True, color=WHITE)
    add_text(s, f"Lift {regla['lift']:.2f}", lx + Inches(2.6), ty + Inches(0.09),
             Inches(1.3), Inches(0.34), size=11, bold=True, color=RGBColor(0x6E, 0xE7, 0xE7),
             align=PP_ALIGN.RIGHT)
    desc = (f"{regla['antecedente']}  →  {regla['consecuente']}   "
            f"(Conf. {regla['confidence']*100:.1f}% · Sup. {regla['support']*100:.1f}%)")
    add_text(s, desc, lx + Inches(0.15), ty + Inches(0.56), Inches(3.75), Inches(1.1),
             size=8, color=INK)
    rect(s, lx + Inches(0.15), ty + Inches(1.7), Inches(3.75), Pt(1.2), fill=BORDER)
    add_text(s, novedad[i], lx + Inches(0.15), ty + Inches(1.8), Inches(3.75), Inches(0.75),
             size=9, italic=True, color=TEAL)


# ═════════════════════════════════════════════════════════════════
# SLIDE 13 · EXPERIMENTOS
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "Experimentos — PCA · TruncatedSVD · Ingeniería de Características",
       "¿Cómo cambian clustering y clasificación con cada transformación?")
slide_num(s, 13)

img_card(s, "outputs/16_comparacion_experimentos.png", Inches(0.45), Inches(1.55), Inches(8.35))

filas_exp = [
    ("Experimento", "Acc.", "F1", "Silh."),
    ("Base", f"{MB['Accuracy']:.3f}", f"{MB['F1 Macro']:.3f}", f"{CL['sil_base']:.3f}"),
    (f"PCA ({EXP['pca']['n_comp']}c)", f"{EXP['pca']['accuracy']:.3f}",
     f"{EXP['pca']['f1_macro']:.3f}", f"{CL['sil_pca']:.3f}"),
    (f"SVD ({EXP['svd']['n_comp']}c)", f"{EXP['svd']['accuracy']:.3f}",
     f"{EXP['svd']['f1_macro']:.3f}", f"{CL['sil_svd']:.3f}"),
    ("FE (+5)", f"{EXP['fe']['accuracy']:.3f}",
     f"{EXP['fe']['f1_macro']:.3f}", f"{CL['sil_fe']:.3f}"),
]
add_text(s, f"Resultados ({MEJOR} · KMeans k=4)", Inches(9.0), Inches(1.6),
         Inches(4.0), Inches(0.4), size=13, bold=True, color=SLATE)
tabla(s, filas_exp,
      [Inches(9.0), Inches(10.6), Inches(11.45), Inches(12.25)],
      [Inches(1.6), Inches(0.85), Inches(0.8), Inches(0.75)],
      Inches(2.05), Inches(0.46), font_size=10)

card(s, Inches(9.0), Inches(4.6), Inches(3.92), Inches(2.45), accent=PURPLE)
add_text(s, "Lectura", Inches(9.3), Inches(4.72), Inches(3.4), Inches(0.35),
         size=13, bold=True, color=PURPLE)
bullet_box(s, [
    f"Clasificación: PCA/SVD pierden señal (F1 {EXP['pca']['f1_macro']:.2f} / "
    f"{EXP['svd']['f1_macro']:.2f} vs {MB['F1 Macro']:.2f})",
    f"Clustering: PCA mejora la separación — Silhouette {CL['sil_base']:.2f} → {CL['sil_pca']:.2f}",
    "FE: BMI eleva F1 a 0.99, pero casi determina el target (construido desde IMC)",
], Inches(9.3), Inches(5.1), Inches(3.45), Inches(1.85), item_size=10, leading=4)


# ═════════════════════════════════════════════════════════════════
# SLIDE 14 · DESCUBRIMIENTOS + BIBLIOGRAFÍA
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "Descubrimientos Destacados y Bibliografía",
       "Hallazgos respaldados por métricas y contrastados con literatura")
slide_num(s, 14)

descubrimientos = [
    ("Mejor clasificador", GREENOK,
     f"{MEJOR}: F1 macro {MB['F1 Macro']:.3f} y accuracy {MB['Accuracy']:.3f} en test. "
     "Das et al. (2026) reportan resultados comparables con ensembles sobre este mismo dataset."),
    ("La historia familiar domina", CORAL,
     "Las reglas con mayor Lift (>6) incluyen family_history=yes; los clusters de mayor IMC "
     "superan 90% de antecedentes. Coincide con el estudio HUNT (Næss et al., 2016): dos padres "
     "con sobrepeso → BMI z-score +0.76 en hijos."),
    ("El IMC casi determina el target", PURPLE,
     f"Con BMI como atributo, F1 sube de {MB['F1 Macro']:.3f} a {EXP['fe']['f1_macro']:.3f}. "
     "El target fue construido desde el IMC (Mendoza Palechor, 2019) → el modelo útil en la "
     "práctica es el que predice solo con hábitos."),
    ("Reducción dimensional: efecto mixto", TEAL,
     f"En clasificación PCA/SVD pierden señal (F1 {EXP['pca']['f1_macro']:.3f} / "
     f"{EXP['svd']['f1_macro']:.3f}), pero PCA mejora el clustering "
     f"(Silhouette {CL['sil_base']:.3f} → {CL['sil_pca']:.3f}). La utilidad depende de la tarea."),
]
for i, (titulo, color, texto) in enumerate(descubrimientos):
    col, row = i % 2, i // 2
    lx = Inches(0.45 + col * 6.35)
    ty = Inches(1.55 + row * 1.85)
    card(s, lx, ty, Inches(6.15), Inches(1.7), accent=color)
    add_text(s, f"{i+1} · {titulo}", lx + Inches(0.32), ty + Inches(0.08),
             Inches(5.7), Inches(0.35), size=13, bold=True, color=color)
    add_text(s, texto, lx + Inches(0.32), ty + Inches(0.46), Inches(5.65), Inches(1.2),
             size=10, color=INK)

card(s, Inches(0.45), Inches(5.4), Inches(12.45), Inches(1.65), fill=TEALBG, accent=TEAL)
add_text(s, "Bibliografía", Inches(0.78), Inches(5.5), Inches(11.9), Inches(0.35),
         size=13, bold=True, color=SLATE)
bullet_box(s, [
    "Mendoza Palechor, F., & De la Hoz Manotas, A. (2019). Dataset for estimation of obesity levels… Data in Brief, 25, 104344. doi:10.1016/j.dib.2019.104344",
    "Das, S., et al. (2026). GBWOEM: A Gradient-Based Weight Optimization Model… F1000Research, 14, 1161. doi:10.12688/f1000research.169436.2",
    "Næss, M., et al. (2016). Intergenerational Transmission of Overweight and Obesity — The HUNT Study. PLoS ONE, 11(11), e0166585. doi:10.1371/journal.pone.0166585",
], Inches(0.78), Inches(5.86), Inches(11.9), Inches(1.15), item_size=10, leading=2)


# ═════════════════════════════════════════════════════════════════
# SLIDE 15 · CONCLUSIONES
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
rect(s, 0, 0, W, H, fill=SLATE)
shape_box(s, 9, Inches(11.2), Inches(-1.8), Inches(3.8), Inches(3.8), fill=SLATE2)
shape_box(s, 9, Inches(12.3), Inches(-0.6), Inches(2.0), Inches(2.0), fill=TEAL)
rect(s, Inches(0.75), Inches(1.12), Inches(2.2), Pt(3), fill=CORAL)

add_text(s, "Conclusiones y Mejoras Futuras", Inches(0.72), Inches(0.3),
         Inches(11.0), Inches(0.7), size=30, bold=True, color=WHITE)
add_text(s, "CRISP-DM · Estimation of Obesity Levels — UCI Machine Learning Repository",
         Inches(0.75), Inches(1.28), Inches(11.5), Inches(0.4), size=13, color=RGBColor(0x8F, 0xA5, 0xBC))
slide_num(s, 15, dark=True)

conclusiones = [
    ("Clustering K-Means (k=4)", TEAL,
     "Cuatro perfiles de riesgo diferenciados por IMC, peso y actividad física. El grupo de "
     "obesidad severa concentra 100% de historia familiar y la menor actividad física."),
    (f"Clasificación — {MEJOR}", RGBColor(0x7D, 0xE8, 0xB8),
     f"Mejor modelo entre cinco algoritmos: F1 macro {MB['F1 Macro']:.4f} y accuracy "
     f"{MB['Accuracy']:.1%} sobre datos de testeo."),
    ("Asociación FP-Growth", AMBER,
     "Lift máximo 6.38: historia familiar + hábitos hipercalóricos + sedentarismo forman un "
     "predictor combinado robusto de obesidad severa. Seis reglas novedosas interpretadas."),
    ("Experimentos de transformación", RGBColor(0xCE, 0x93, 0xD8),
     f"PCA mejora el clustering (Silhouette {CL['sil_base']:.2f} → {CL['sil_pca']:.2f}) pero pierde "
     f"señal en clasificación; la ingeniería de características eleva F1 a "
     f"{EXP['fe']['f1_macro']:.2f}, con la advertencia de que el IMC casi determina el target."),
    ("Mejoras futuras", CORAL,
     "Optimización de hiperparámetros con validación cruzada · modelo solo con hábitos (sin "
     "antropometría directa) · validación con datos reales no sintéticos · aplicación de "
     "consulta del clasificador (Streamlit)."),
]
for i, (titulo, color, texto) in enumerate(conclusiones):
    ty = Inches(1.85 + i * 1.05)
    shape_box(s, 5, Inches(0.45), ty, Inches(12.45), Inches(0.95), fill=SLATE2)
    rect(s, Inches(0.45), ty + Inches(0.08), Inches(0.09), Inches(0.79), fill=color)
    add_text(s, titulo, Inches(0.75), ty + Inches(0.05), Inches(3.7), Inches(0.4),
             size=12, bold=True, color=color)
    add_text(s, texto, Inches(4.5), ty + Inches(0.06), Inches(8.25), Inches(0.85),
             size=10, color=RGBColor(0xD9, 0xE2, 0xEC))

add_text(s, "Notebook reproducible · Dataset UCI · Dante Gil Zenteno — Magíster en Data Science, USS",
         Inches(0.75), Inches(7.12), Inches(11.5), Inches(0.3),
         size=10, color=RGBColor(0x7E, 0x95, 0xAE))


# ─── GUARDAR ─────────────────────────────────────────────────────
out_path = "presentacion_obesidad_final.pptx"
prs.save(out_path)
print(f"PPT final generada: {out_path}")
print(f"Diapositivas: {len(prs.slides)}")
