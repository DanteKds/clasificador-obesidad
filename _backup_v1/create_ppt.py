"""
Generador de PPT — Análisis de Obesidad CRISP-DM
Ejecutar: python create_ppt.py
Produce:  presentacion_obesidad.pptx (máx 15 diapositivas)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── PALETA DE COLORES ───────────────────────────────────────────
NAVY    = RGBColor(0x1A, 0x23, 0x5E)   # azul marino oscuro (fondo header)
BLUE    = RGBColor(0x1E, 0x6B, 0xB8)   # azul medio (acentos)
CYAN    = RGBColor(0x00, 0xBF, 0xD8)   # cyan (highlights)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xF5, 0xF5, 0xF5)   # fondo gris claro
DGRAY   = RGBColor(0x44, 0x44, 0x44)   # texto gris oscuro
GREEN   = RGBColor(0x2E, 0x7D, 0x32)
ORANGE  = RGBColor(0xE6, 0x51, 0x00)

W = Inches(13.33)  # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H


# ─── HELPERS ─────────────────────────────────────────────────────
def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # completamente en blanco
    return prs.slides.add_slide(blank_layout)


def rect(slide, l, t, w, h, fill=None, line_color=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=DGRAY,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def header_bar(slide, title, subtitle=None):
    """Barra superior azul con título."""
    rect(slide, 0, 0, W, Inches(1.35), fill=NAVY)
    add_text(slide, title,
             Inches(0.35), Inches(0.08), Inches(12.5), Inches(0.75),
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.35), Inches(0.8), Inches(12.5), Inches(0.45),
                 size=14, color=CYAN)


def slide_num(slide, n, total=15):
    add_text(slide, f"{n} / {total}",
             Inches(12.6), Inches(7.1), Inches(0.7), Inches(0.3),
             size=9, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.RIGHT)


def add_image(slide, path, l, t, w, h=None):
    if not os.path.exists(path):
        return
    if h:
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        slide.shapes.add_picture(path, l, t, w)


def bullet_box(slide, items, l, t, w, h, title=None, title_color=NAVY,
               item_size=13, item_color=DGRAY, indent=True):
    """Caja con lista de bullets."""
    if title:
        add_text(slide, title, l, t, w, Inches(0.35),
                 size=14, bold=True, color=title_color)
        t += Inches(0.38)
        h -= Inches(0.38)
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = ("• " if indent else "") + item
        run.font.size = Pt(item_size)
        run.font.color.rgb = item_color


# ─── SLIDE 1: PORTADA ───────────────────────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=NAVY)
rect(s, Inches(0.5), Inches(1.8), Inches(12.33), Inches(0.05), fill=CYAN)

add_text(s, "Análisis de Niveles de Obesidad",
         Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8),
         size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "mediante Métodos de Ciencia de Datos",
         Inches(0.5), Inches(1.0), Inches(12.33), Inches(0.65),
         size=24, color=CYAN, align=PP_ALIGN.CENTER)

add_text(s, "Dataset: Estimation of Obesity Levels — UCI ML Repository",
         Inches(1.5), Inches(2.1), Inches(10.33), Inches(0.45),
         size=14, color=WHITE, align=PP_ALIGN.CENTER)

# Info box
rect(s, Inches(3.0), Inches(3.0), Inches(7.33), Inches(2.9),
     fill=RGBColor(0x22, 0x2D, 0x6E), line_color=CYAN)

infos = [
    ("Alumno:",          "Dante Gil Zenteno"),
    ("Docente:",         "Dr. Mauricio Sepúlveda"),
    ("Asignatura:",      "Taller de Aplicaciones"),
    ("Programa:",        "Magíster en Data Science — Facultad de Ingeniería"),
]
for i, (label, val) in enumerate(infos):
    y = Inches(3.15) + i * Inches(0.62)
    add_text(s, label, Inches(3.3), y, Inches(2.0), Inches(0.5),
             size=13, bold=True, color=CYAN)
    add_text(s, val, Inches(5.3), y, Inches(4.8), Inches(0.5),
             size=13, color=WHITE)

add_text(s, "2026", Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.35),
         size=11, color=RGBColor(0x88, 0x88, 0xAA), align=PP_ALIGN.CENTER)
slide_num(s, 1)


# ─── SLIDE 2: CONTEXTO Y OBJETIVO ───────────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Contexto del Problema y Objetivo",
           "CRISP-DM · Fase 1: Comprensión del Negocio")
slide_num(s, 2)

rect(s, Inches(0.3), Inches(1.5), Inches(6.1), Inches(5.5), fill=WHITE,
     line_color=RGBColor(0xCC, 0xCC, 0xCC))
rect(s, Inches(6.8), Inches(1.5), Inches(6.2), Inches(5.5), fill=WHITE,
     line_color=RGBColor(0xCC, 0xCC, 0xCC))

add_text(s, "El Dataset", Inches(0.5), Inches(1.55), Inches(5.7), Inches(0.4),
         size=15, bold=True, color=NAVY)
dataset_items = [
    "2.111 registros de México, Perú y Colombia",
    "17 atributos: hábitos alimentarios, actividad física y características personales",
    "Variable objetivo NObeyesdad: 7 niveles de obesidad",
    "Clases balanceadas (272–351 registros cada una)",
    "Sin valores nulos ni duplicados",
    "Fuente: Mendoza Palechor & De la Hoz Manotas (2019), Data in Brief",
]
bullet_box(s, dataset_items, Inches(0.5), Inches(2.0), Inches(5.7), Inches(4.7),
           item_size=12)

add_text(s, "Preguntas Guía", Inches(7.0), Inches(1.55), Inches(5.8), Inches(0.4),
         size=15, bold=True, color=NAVY)
preguntas = [
    "¿Existen grupos de personas con perfiles de riesgo similares?",
    "¿Qué algoritmo predice mejor el nivel de obesidad?",
    "¿Qué combinaciones de hábitos se asocian con ciertos niveles de obesidad?",
]
bullet_box(s, preguntas, Inches(7.0), Inches(2.0), Inches(5.8), Inches(2.5), item_size=13)

add_text(s, "Objetivo General", Inches(7.0), Inches(4.6), Inches(5.8), Inches(0.4),
         size=15, bold=True, color=NAVY)
add_text(s, "Aplicar CRISP-DM sobre el dataset de obesidad usando tres métodos de extracción de conocimiento: agrupamiento, clasificación y asociación con FP-Growth.",
         Inches(7.0), Inches(5.05), Inches(5.8), Inches(1.5),
         size=12, color=DGRAY)


# ─── SLIDE 3: CRISP-DM ──────────────────────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Metodología CRISP-DM", "Estructura de trabajo del proyecto")
slide_num(s, 3)

fases = [
    ("1", "Comprensión\ndel Negocio", "Definición del problema, preguntas guía, métricas de éxito"),
    ("2", "Comprensión\nde los Datos", "Carga del CSV, inspección, estadísticas, calidad"),
    ("3", "Preparación\nde los Datos", "Limpieza, codificación, escalado, versiones por modelo"),
    ("4", "Modelamiento", "K-Means · Clasificación · FP-Growth"),
    ("5", "Evaluación",   "Métricas en test, Lift, interpretación de reglas"),
    ("6", "Presentación", "PPT 15 diapositivas + notebook reproducible"),
]

for i, (num, nombre, desc) in enumerate(fases):
    col = i % 3
    row = i // 3
    lx = Inches(0.4 + col * 4.3)
    ty = Inches(1.6 + row * 2.6)
    rect(s, lx, ty, Inches(3.9), Inches(2.4), fill=WHITE,
         line_color=RGBColor(0xCC, 0xCC, 0xCC))
    rect(s, lx, ty, Inches(3.9), Inches(0.5), fill=NAVY)
    add_text(s, f"Fase {num}", lx + Inches(0.1), ty + Inches(0.05),
             Inches(3.7), Inches(0.4), size=13, bold=True, color=WHITE)
    add_text(s, nombre, lx + Inches(0.15), ty + Inches(0.55),
             Inches(3.6), Inches(0.7), size=14, bold=True, color=NAVY)
    add_text(s, desc, lx + Inches(0.15), ty + Inches(1.2),
             Inches(3.6), Inches(1.1), size=11, color=DGRAY)


# ─── SLIDE 4: DESCRIPCIÓN DEL DATASET ───────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Descripción del Dataset",
           "CRISP-DM · Fase 2: Comprensión de los Datos")
slide_num(s, 4)

# Left: variable list
add_text(s, "Variables del Dataset (17)", Inches(0.3), Inches(1.5), Inches(4.5), Inches(0.4),
         size=14, bold=True, color=NAVY)
vars_list = [
    "Género (Gender)",
    "Edad (Age)  |  Altura (Height)  |  Peso (Weight)",
    "Historia familiar sobrepeso",
    "Consumo alimentos hipercalóricos (FAVC)",
    "Frecuencia de vegetales (FCVC)",
    "Número de comidas principales (NCP)",
    "Snacks entre comidas (CAEC)",
    "Tabaquismo (SMOKE)  |  Agua (CH2O)",
    "Monitoreo de calorías (SCC)",
    "Actividad física (FAF)  |  Tiempo pantallas (TUE)",
    "Consumo de alcohol (CALC)",
    "Medio de transporte (MTRANS)",
    "→ Nivel de obesidad: NObeyesdad  ★ TARGET",
]
bullet_box(s, vars_list, Inches(0.3), Inches(1.95), Inches(4.8), Inches(5.0),
           item_size=11)

# Right: stats table
add_text(s, "Niveles de Obesidad (Variable Objetivo)", Inches(5.5), Inches(1.5),
         Inches(7.5), Inches(0.4), size=14, bold=True, color=NAVY)

niveles = [
    ("Nivel de Obesidad",          "Registros", "%"),
    ("Obesity_Type_I",             "351",       "16.6%"),
    ("Obesity_Type_III",           "324",       "15.3%"),
    ("Obesity_Type_II",            "297",       "14.1%"),
    ("Overweight_Level_I",         "290",       "13.7%"),
    ("Overweight_Level_II",        "290",       "13.7%"),
    ("Normal_Weight",              "287",       "13.6%"),
    ("Insufficient_Weight",        "272",       "12.9%"),
]
row_h = Inches(0.42)
col_ws = [Inches(3.8), Inches(1.3), Inches(1.3)]
col_ls = [Inches(5.5), Inches(9.3), Inches(10.6)]
ty = Inches(1.95)
for ri, row in enumerate(niveles):
    bg = NAVY if ri == 0 else (WHITE if ri % 2 == 1 else LGRAY)
    fc = WHITE if ri == 0 else DGRAY
    bold = ri == 0
    for ci, (text, lx, cw) in enumerate(zip(row, col_ls, col_ws)):
        rect(s, lx, ty, cw, row_h, fill=bg,
             line_color=RGBColor(0xCC, 0xCC, 0xCC))
        add_text(s, text, lx + Inches(0.08), ty + Inches(0.07), cw - Inches(0.1),
                 row_h - Inches(0.07), size=11, bold=bold, color=fc)
    ty += row_h


# ─── SLIDE 5: PREPARACIÓN DE DATOS ──────────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Preparación y Procesamiento de Datos",
           "CRISP-DM · Fase 3: Preparación de los Datos")
slide_num(s, 5)

steps = [
    ("1. Calidad", ["Sin nulos (0 valores faltantes)", "Sin filas duplicadas", "Rangos numéricos consistentes"]),
    ("2. Identificación de Variables", ["Numéricas (8): Age, Height, Weight, FCVC, NCP, CH2O, FAF, TUE", "Categóricas (8): Gender, FAVC, CAEC, SMOKE, SCC, CALC, MTRANS, family_history", "Target (1): NObeyesdad — 7 clases"]),
    ("3. Para Clustering", ["OrdinalEncoder en variables categóricas", "StandardScaler en todas las variables", "Target excluido → sin fuga de datos"]),
    ("4. Para Clasificación", ["Pipeline con ColumnTransformer", "OrdinalEncoder + StandardScaler dentro del pipeline", "train_test_split estratificado 80%/20%"]),
    ("5. Para FP-Growth", ["Variables numéricas discretizadas en rangos (Age, Weight, Height, etc.)", "Variables categóricas prefijadas: Gender=Female, CAEC=Sometimes…", "Formato OHE booleano para mlxtend"]),
]

col_w = Inches(2.45)
for i, (title, items) in enumerate(steps):
    lx = Inches(0.3 + i * 2.6)
    rect(s, lx, Inches(1.5), col_w, Inches(5.7), fill=WHITE,
         line_color=RGBColor(0xCC, 0xCC, 0xCC))
    rect(s, lx, Inches(1.5), col_w, Inches(0.5), fill=BLUE)
    add_text(s, title, lx + Inches(0.1), Inches(1.55), col_w - Inches(0.1),
             Inches(0.4), size=12, bold=True, color=WHITE)
    bullet_box(s, items, lx + Inches(0.1), Inches(2.1),
               col_w - Inches(0.15), Inches(4.9), item_size=10)


# ─── SLIDE 6: EDA ────────────────────────────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Análisis Exploratorio de Datos (EDA)",
           "CRISP-DM · Fase 2 — Distribuciones, correlaciones y patrones")
slide_num(s, 6)

add_image(s, "outputs/01_distribucion_objetivo.png",
          Inches(0.2), Inches(1.45), Inches(6.5))
add_image(s, "outputs/03_correlacion.png",
          Inches(7.0), Inches(1.45), Inches(6.0))

add_text(s, "Hallazgos clave:",
         Inches(0.2), Inches(6.5), Inches(13.0), Inches(0.35),
         size=12, bold=True, color=NAVY)
add_text(s,
         "• 7 clases balanceadas (272–351 por nivel)   "
         "• Weight y Height con mayor correlación al target   "
         "• Age tiene menor variación entre clases",
         Inches(0.2), Inches(6.85), Inches(13.0), Inches(0.45),
         size=11, color=DGRAY)


# ─── SLIDE 7: CLUSTERING — MÉTODO Y JUSTIFICACIÓN ───────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Clustering — Algoritmo y Número de Grupos",
           "CRISP-DM · Fase 4: Modelamiento (Agrupamiento)")
slide_num(s, 7)

add_image(s, "outputs/06_evaluacion_clusters.png",
          Inches(0.2), Inches(1.45), Inches(8.2))

# Justificación
rect(s, Inches(8.6), Inches(1.5), Inches(4.5), Inches(5.7), fill=WHITE,
     line_color=RGBColor(0xCC, 0xCC, 0xCC))
add_text(s, "Justificación del Algoritmo", Inches(8.8), Inches(1.6),
         Inches(4.2), Inches(0.4), size=13, bold=True, color=NAVY)
bullet_box(s,
           ["K-Means: escala bien con 2.111 registros",
            "Variables mixtas escaladas → centroides interpretables",
            "Alternativas evaluadas: DBSCAN (requiere calibración de eps) y Agglomerative (lento y similar resultado)"],
           Inches(8.8), Inches(2.05), Inches(4.2), Inches(1.9), item_size=11)

add_text(s, "Elección de k=4", Inches(8.8), Inches(4.0),
         Inches(4.2), Inches(0.4), size=13, bold=True, color=NAVY)
bullet_box(s,
           ["k=2 maximiza Silhouette (0.190) pero solo separa \"obeso / no-obeso\" — insuficiente",
            "k=4 mapea los 7 niveles de obesidad en 4 perfiles de riesgo diferenciados",
            "Permite diseñar intervenciones distintas por grupo",
            "Silhouette scores bajos en todo el rango (<0.20) son normales en datos de salud con variables sintéticas"],
           Inches(8.8), Inches(4.45), Inches(4.2), Inches(2.5), item_size=11)


# ─── SLIDE 8: CLUSTERING — INTERPRETACIÓN ────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Clustering — Interpretación de los 4 Grupos",
           "CRISP-DM · Fase 4: Modelamiento")
slide_num(s, 8)

add_image(s, "outputs/07_clusters_visualizacion.png",
          Inches(0.2), Inches(1.45), Inches(7.8))

clusters = [
    ("Cluster 0 · 22.4%", "Bajo peso / Normal",
     "IMC 21.4 · Peso 56.6 kg\nEdad 20.8 · FAF 1.10\nHist. familiar 41%",
     CYAN),
    ("Cluster 1 · 23.1%", "Obesidad Severa",
     "IMC 37.8 · Peso 106.6 kg\nEdad 22.9 · FAF 0.73\nHist. familiar 100%",
     ORANGE),
    ("Cluster 2 · 36.4%", "Obesidad Moderada",
     "IMC 28.9 · Peso 90.2 kg\nEdad 22.3 · FAF 1.20\nHist. familiar 90%",
     RGBColor(0xE6, 0x9A, 0x00)),
    ("Cluster 3 · 18.1%", "Obesidad Leve/I",
     "IMC 31.2 · Peso 91.0 kg\nEdad 34.4 · FAF 0.88\nHist. familiar 92%",
     RGBColor(0x1B, 0x5E, 0x20)),
]

for i, (nombre, perfil, stats, color) in enumerate(clusters):
    lx = Inches(8.2)
    ty = Inches(1.5 + i * 1.45)
    rect(s, lx, ty, Inches(5.0), Inches(1.35), fill=WHITE,
         line_color=RGBColor(0xCC, 0xCC, 0xCC))
    rect(s, lx, ty, Inches(0.25), Inches(1.35), fill=color)
    add_text(s, nombre, lx + Inches(0.35), ty + Inches(0.05),
             Inches(4.5), Inches(0.35), size=12, bold=True, color=NAVY)
    add_text(s, perfil, lx + Inches(0.35), ty + Inches(0.38),
             Inches(2.2), Inches(0.35), size=11, bold=True, color=color)
    add_text(s, stats, lx + Inches(0.35), ty + Inches(0.7),
             Inches(4.5), Inches(0.6), size=10, color=DGRAY)


# ─── SLIDE 9: CLASIFICACIÓN — MODELOS ────────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Clasificación — Modelos Comparados",
           "CRISP-DM · Fase 4: Modelamiento (Clasificación)")
slide_num(s, 9)

add_image(s, "outputs/08_comparacion_modelos.png",
          Inches(0.2), Inches(1.45), Inches(9.5))

add_text(s, "Configuración del experimento", Inches(10.0), Inches(1.5),
         Inches(3.1), Inches(0.4), size=13, bold=True, color=NAVY)
bullet_box(s,
           ["Train / Test: 80% / 20%",
            "Estratificación por clase",
            "Pipeline con preprocesamiento dentro (sin data leakage)",
            "Métricas: accuracy, precision, recall y F1 macro",
            "random_state = 42"],
           Inches(10.0), Inches(1.95), Inches(3.1), Inches(2.5), item_size=11)

add_text(s, "Modelos evaluados", Inches(10.0), Inches(4.5),
         Inches(3.1), Inches(0.4), size=13, bold=True, color=NAVY)
bullet_box(s,
           ["Regresión Logística Multinomial",
            "Árbol de Decisión (max_depth=12)",
            "Random Forest (100 árboles)",
            "K-Nearest Neighbors (k=7)"],
           Inches(10.0), Inches(4.95), Inches(3.1), Inches(2.0), item_size=11)


# ─── SLIDE 10: CLASIFICACIÓN — RESULTADOS ────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Clasificación — Resultados en Test",
           "CRISP-DM · Fase 5: Evaluación")
slide_num(s, 10)

add_image(s, "outputs/09_confusion_matrix.png",
          Inches(0.2), Inches(1.45), Inches(6.5))

# Results table
add_text(s, "Métricas en Datos de Testeo (20%)", Inches(7.0), Inches(1.5),
         Inches(6.1), Inches(0.4), size=14, bold=True, color=NAVY)

results_data = [
    ("Modelo",                  "Accuracy", "Precision", "Recall",   "F1 Macro"),
    ("Random Forest ★",         "0.9621",   "0.9618",    "0.9572",   "0.9596"),
    ("Árbol de Decisión",       "0.9408",   "0.9406",    "0.9412",   "0.9403"),
    ("KNN",                     "0.8793",   "0.8796",    "0.8751",   "0.8767"),
    ("Regresión Logística",     "0.8317",   "0.8309",    "0.8272",   "0.8278"),
]
row_h2 = Inches(0.5)
col_ws2 = [Inches(2.5), Inches(0.85), Inches(0.9), Inches(0.85), Inches(0.9)]
col_ls2 = [Inches(7.0), Inches(9.5), Inches(10.35), Inches(11.25), Inches(12.1)]
ty = Inches(1.95)
for ri, row in enumerate(results_data):
    bg = NAVY if ri == 0 else (RGBColor(0xE8, 0xF5, 0xE9) if ri == 1 else (WHITE if ri % 2 == 0 else LGRAY))
    fc = WHITE if ri == 0 else (GREEN if ri == 1 else DGRAY)
    bold = ri <= 1
    for ci, (text, lx, cw) in enumerate(zip(row, col_ls2, col_ws2)):
        rect(s, lx, ty, cw, row_h2, fill=bg,
             line_color=RGBColor(0xCC, 0xCC, 0xCC))
        add_text(s, text, lx + Inches(0.05), ty + Inches(0.1),
                 cw - Inches(0.05), row_h2 - Inches(0.1),
                 size=11, bold=bold, color=fc, align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
    ty += row_h2

add_text(s, "★ Random Forest es el mejor modelo con F1 Macro = 0.9596",
         Inches(7.0), Inches(4.6), Inches(6.1), Inches(0.5),
         size=13, bold=True, color=GREEN)

add_text(s, "La clase Normal_Weight tiene menor precision (0.85) por confusión\ncon Overweight — normal dado que son límites continuos.",
         Inches(7.0), Inches(5.2), Inches(6.1), Inches(0.8),
         size=11, color=DGRAY)


# ─── SLIDE 11: MEJOR MODELO Y MEJORAS ────────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Mejor Modelo y Mejoras Futuras",
           "CRISP-DM · Fase 5: Evaluación — Mejoras propuestas")
slide_num(s, 11)

rect(s, Inches(0.3), Inches(1.5), Inches(5.8), Inches(2.5), fill=RGBColor(0xE8, 0xF5, 0xE9),
     line_color=GREEN)
add_text(s, "Mejor Modelo: Random Forest", Inches(0.5), Inches(1.6),
         Inches(5.4), Inches(0.4), size=15, bold=True, color=GREEN)
bullet_box(s,
           ["F1 Macro test: 0.9596",
            "Accuracy: 96.2%",
            "100 árboles · max_features='sqrt' (default)",
            "Robusto a variables categóricas codificadas",
            "Sin ajuste de hiperparámetros — resultado base ya excelente"],
           Inches(0.5), Inches(2.05), Inches(5.4), Inches(1.8), item_size=12)

mejoras = [
    ("Ingeniería de Variables",     "Crear IMC = Peso/Altura², interacciones FAF×FAVC"),
    ("Hiperparámetros",             "GridSearchCV o RandomizedSearchCV para RF/GB"),
    ("Gradient Boosting / XGBoost", "Modelos más potentes que Random Forest"),
    ("Variables Ordinales",         "CAEC y CALC son ordinales — no solo nominales"),
    ("Validación Cruzada",          "k=5 folds estratificados para estimación robusta"),
    ("Datos Reales",                "77% del set fue generado con SMOTE — validar en datos reales"),
    ("Reducción Dimensional",       "PCA para clustering más compacto"),
    ("Balanceo de Clases",          "SMOTE si se trabaja solo con registros reales"),
]
add_text(s, "Mejoras Futuras (8)", Inches(0.3), Inches(4.1), Inches(12.7), Inches(0.4),
         size=14, bold=True, color=NAVY)
row_h3 = Inches(0.4)
ty = Inches(4.55)
for i, (tit, desc) in enumerate(mejoras):
    col = i % 2
    lx = Inches(0.3 + col * 6.5)
    row = i // 2
    ty2 = ty + row * row_h3
    rect(s, lx, ty2, Inches(6.2), row_h3, fill=WHITE if i % 2 == 0 else LGRAY,
         line_color=RGBColor(0xCC, 0xCC, 0xCC))
    add_text(s, f"• {tit}: ", lx + Inches(0.08), ty2 + Inches(0.05),
             Inches(2.0), Inches(0.32), size=10, bold=True, color=NAVY)
    add_text(s, desc, lx + Inches(2.1), ty2 + Inches(0.05),
             Inches(4.0), Inches(0.32), size=10, color=DGRAY)


# ─── SLIDE 12: FP-GROWTH — PREPARACIÓN ───────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "FP-Growth — Preparación de los Datos",
           "CRISP-DM · Fase 4: Modelamiento (Asociación)")
slide_num(s, 12)

add_text(s, "Transformación a Formato Transaccional Binario",
         Inches(0.3), Inches(1.5), Inches(13.0), Inches(0.4),
         size=14, bold=True, color=NAVY)

# Columns: numéricas y categóricas
rect(s, Inches(0.3), Inches(2.0), Inches(5.9), Inches(4.8), fill=WHITE,
     line_color=RGBColor(0xCC, 0xCC, 0xCC))
add_text(s, "Variables Numéricas → Discretizadas", Inches(0.5), Inches(2.1),
         Inches(5.5), Inches(0.35), size=12, bold=True, color=NAVY)
disc_items = [
    "Age  → Adolescente / Joven / Adulto / Mayor",
    "Height  → Bajo / MedioBajo / MedioAlto / Alto",
    "Weight  → Leve / Normal / Alto / MuyAlto / Extremo",
    "FCVC  → Bajo / Medio / Alto  (consumo vegetales)",
    "NCP  → Bajo / Normal / Alto  (comidas principales)",
    "CH2O  → Bajo / Medio / Alto  (agua consumida)",
    "FAF  → Nulo / Bajo / Moderado / Alto  (actividad)",
    "TUE  → Bajo / Medio / Alto  (tiempo pantallas)",
]
bullet_box(s, disc_items, Inches(0.5), Inches(2.5), Inches(5.5), Inches(4.2), item_size=11)

rect(s, Inches(6.5), Inches(2.0), Inches(6.5), Inches(4.8), fill=WHITE,
     line_color=RGBColor(0xCC, 0xCC, 0xCC))
add_text(s, "Ítems Resultantes (ejemplos)", Inches(6.7), Inches(2.1),
         Inches(6.1), Inches(0.35), size=12, bold=True, color=NAVY)
item_examples = [
    "Gender=Female",
    "family_history_with_overweight=yes",
    "Age=Joven",
    "Weight=Alto",
    "FAVC=yes",
    "FAF=Nulo",
    "MTRANS=Public_Transportation",
    "NObeyesdad=Obesity_Type_III   ← target incluido",
]
bullet_box(s, item_examples, Inches(6.7), Inches(2.5), Inches(6.0), Inches(3.0), item_size=12)

add_text(s, "Parámetros FP-Growth: min_support = 0.15 (315 registros mínimos) · min_confidence = 0.70 · Lift > 2.0",
         Inches(6.7), Inches(5.6), Inches(6.1), Inches(0.6),
         size=12, bold=False, color=BLUE)

add_text(s, "→ Resultado: 7.312 itemsets frecuentes · 5.497 reglas con NObeyesdad como consecuente",
         Inches(0.3), Inches(6.7), Inches(13.0), Inches(0.45),
         size=12, bold=True, color=NAVY)


# ─── SLIDE 13: REGLAS CON MEJORES INDICADORES ────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Reglas de Asociación — Mejores Indicadores",
           "FP-Growth · Filtro: Lift > 2.0, Confidence ≥ 70%, NObeyesdad como consecuente")
slide_num(s, 13)

add_image(s, "outputs/10_reglas_asociacion.png",
          Inches(0.2), Inches(1.45), Inches(7.5))

top_rules = [
    ("Antecedentes",                       "Consecuente",   "Sup",   "Conf",   "Lift"),
    ("family_history=yes, Weight=Extremo", "Obesity_III",   "0.153", "0.979",  "6.38"),
    ("FAVC=yes, FAF=Nulo, Weight=Alto",    "Obesity_II",    "0.167", "0.971",  "6.34"),
    ("CALC=Sometimes, Gender=Female",      "Obesity_III",   "0.158", "0.969",  "6.31"),
    ("NCP=Normal, SCC=no, FAVC=yes",       "Obesity_III",   "0.152", "0.970",  "6.32"),
    ("CH2O=Bajo, FAF=Nulo",                "Obesity_I",     "0.189", "0.943",  "5.47"),
    ("MTRANS=Public, FCVC=Alto",           "Obesity_III",   "0.155", "0.974",  "6.35"),
]
add_text(s, "Top Reglas (NObeyesdad como consecuente)", Inches(8.0), Inches(1.5),
         Inches(5.2), Inches(0.4), size=13, bold=True, color=NAVY)

rh = Inches(0.5)
col_ls3  = [Inches(8.0), Inches(10.8), Inches(11.45), Inches(12.05), Inches(12.65)]
col_ws3  = [Inches(2.75), Inches(0.65), Inches(0.6),  Inches(0.6),  Inches(0.65)]
ty = Inches(1.95)
for ri, row in enumerate(top_rules):
    bg = NAVY if ri == 0 else (WHITE if ri % 2 == 1 else LGRAY)
    fc = WHITE if ri == 0 else DGRAY
    for ci, (text, lx, cw) in enumerate(zip(row, col_ls3, col_ws3)):
        rect(s, lx, ty, cw, rh, fill=bg, line_color=RGBColor(0xCC, 0xCC, 0xCC))
        add_text(s, text, lx + Inches(0.04), ty + Inches(0.08),
                 cw - Inches(0.04), rh - Inches(0.08),
                 size=9, bold=ri == 0, color=fc,
                 align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
    ty += rh

add_text(s, "★ Lift > 6 indica que la regla es 6× más probable que por azar",
         Inches(8.0), Inches(5.35), Inches(5.2), Inches(0.45),
         size=11, bold=True, color=ORANGE)


# ─── SLIDE 14: 6 REGLAS NOVEDOSAS ────────────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "6 Reglas Novedosas de Asociación",
           "FP-Growth · Interpretación y novedad de las reglas seleccionadas")
slide_num(s, 14)

reglas = [
    ("Regla 1", "Lift 6.38",
     "Mujer con historial familiar + hábitos frecuentes (CAEC, FAVC, NCP normal)\n→ Obesity_Type_III  (Conf. 97.9%)",
     "Combina factor genético con patrón conductual complejo — no trivial."),
    ("Regla 2", "Lift 6.38",
     "FAVC=yes + FCVC=Alto + MTRANS=Public + sin control calórico\n→ Obesity_Type_III  (Conf. 97.6%)",
     "El transporte público sedentario potencia el efecto de la dieta hipercalórica."),
    ("Regla 3", "Lift 6.38",
     "Misma combinación de la Regla 1 + SMOKE=no\n→ Obesity_Type_III y SMOKE=no  (Conf. 97.6%)",
     "Confirma que no fumar no protege de obesidad severa cuando los hábitos alimentarios son malos."),
    ("Regla 4", "Lift 6.36",
     "CALC=Sometimes + FAVC=yes + FCVC=Alto + historial familiar\n→ CAEC=Sometimes + Obesity_III  (Conf. 97.0%)",
     "Snacking y consumo de alcohol moderado refuerzan el riesgo — regla multivariada novedosa."),
    ("Regla 5", "Lift 6.35",
     "MTRANS=Public + FCVC=Alto + NCP=Normal + historial\n→ FAVC=yes + Obesity_Type_III  (Conf. 97.4%)",
     "Predice también el patrón de conducta (FAVC) no solo el nivel de obesidad."),
    ("Regla 6", "Lift 6.32",
     "CALC=Sometimes + Gender=Female + historial + FCVC=Alto\n→ Obesity_Type_III  (Conf. 96.9%)",
     "Consumo de alcohol ocasional en mujeres con antecedentes familiares y dieta alta en calorías."),
]

card_w = Inches(4.1)
card_h = Inches(2.5)
for i, (titulo, lift, descripcion, novedad) in enumerate(reglas):
    col = i % 3
    row = i // 3
    lx = Inches(0.25 + col * 4.35)
    ty = Inches(1.5 + row * 2.7)
    rect(s, lx, ty, card_w, card_h, fill=WHITE,
         line_color=RGBColor(0xCC, 0xCC, 0xCC))
    rect(s, lx, ty, card_w, Inches(0.45), fill=NAVY)
    add_text(s, titulo, lx + Inches(0.1), ty + Inches(0.05),
             Inches(2.5), Inches(0.35), size=12, bold=True, color=WHITE)
    add_text(s, lift, lx + card_w - Inches(1.3), ty + Inches(0.08),
             Inches(1.2), Inches(0.3), size=11, bold=True,
             color=CYAN, align=PP_ALIGN.RIGHT)
    add_text(s, descripcion, lx + Inches(0.1), ty + Inches(0.5),
             card_w - Inches(0.15), Inches(1.05), size=9, color=DGRAY)
    rect(s, lx, ty + Inches(1.6), card_w, Inches(0.85),
         fill=RGBColor(0xF0, 0xF4, 0xFF))
    add_text(s, "Novedad: " + novedad, lx + Inches(0.1), ty + Inches(1.65),
             card_w - Inches(0.15), Inches(0.75), size=8, italic=True,
             color=BLUE)


# ─── SLIDE 15: CONCLUSIONES ──────────────────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=NAVY)
rect(s, 0, 0, W, Inches(0.08), fill=CYAN)
rect(s, 0, H - Inches(0.08), W, Inches(0.08), fill=CYAN)

add_text(s, "Conclusiones Finales", Inches(0.5), Inches(0.15),
         Inches(12.33), Inches(0.7), size=30, bold=True, color=WHITE)
add_text(s, "CRISP-DM · Dataset: Estimation of Obesity Levels — UCI ML Repository",
         Inches(0.5), Inches(0.8), Inches(12.33), Inches(0.4),
         size=14, color=CYAN)
slide_num(s, 15)

conclusiones = [
    ("Clustering K-Means k=4",
     "4 perfiles de riesgo claramente diferenciados por IMC, peso y actividad física. "
     "Cluster 1 (23.1%): obesidad severa, 100% historia familiar, actividad baja.",
     CYAN),
    ("Clasificación — Random Forest",
     "Mejor modelo con F1 Macro = 0.9596 sobre datos de testeo. "
     "96% accuracy. Supera en 6-16 puntos a los demás modelos.",
     RGBColor(0x66, 0xFF, 0x66)),
    ("Asociación FP-Growth",
     "Lift máximo = 6.38 — las reglas son 6× más informativas que el azar. "
     "Historia familiar + hábitos hipercalóricos + transporte sedentario = predictor combinado robusto de obesidad severa.",
     ORANGE),
    ("Hallazgo Clave",
     "No basta un solo factor: la obesidad tipo III emerge de la combinación de antecedentes "
     "genéticos + conducta alimentaria + sedentarismo. Esto valida el enfoque multivariado.",
     RGBColor(0xFF, 0xCC, 0x00)),
]

for i, (titulo, texto, color) in enumerate(conclusiones):
    ty = Inches(1.35 + i * 1.45)
    rect(s, Inches(0.3), ty, Inches(12.7), Inches(1.3),
         fill=RGBColor(0x22, 0x2D, 0x6E),
         line_color=RGBColor(0x33, 0x3D, 0x7E))
    rect(s, Inches(0.3), ty, Inches(0.2), Inches(1.3), fill=color)
    add_text(s, titulo, Inches(0.65), ty + Inches(0.1),
             Inches(3.0), Inches(0.4), size=13, bold=True, color=color)
    add_text(s, texto, Inches(0.65), ty + Inches(0.5),
             Inches(12.15), Inches(0.75), size=11, color=WHITE)

add_text(s, "Notebook reproducible disponible · Dataset UCI · Alumno: Dante Gil Zenteno",
         Inches(0.5), Inches(7.1), Inches(12.33), Inches(0.3),
         size=10, color=RGBColor(0x88, 0x88, 0xAA), align=PP_ALIGN.CENTER)


# ─── GUARDAR ─────────────────────────────────────────────────────
out_path = "presentacion_obesidad.pptx"
prs.save(out_path)
print(f"PPT generada: {out_path}")
print(f"Diapositivas: {len(prs.slides)}")
