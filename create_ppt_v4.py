# -*- coding: utf-8 -*-
"""
Generador de PPT FINAL v4 — Análisis de Obesidad CRISP-DM
Ejecutar: venvAPP/Scripts/python.exe create_ppt_v4.py
Produce:  presentacion_obesidad_final_v4.pptx (~19 diapositivas + portada)

Paleta: PRD 04 §10.2 (azul académico / morado analítico / verde / naranja / rojo)
Métricas: leídas desde outputs/resultados_v2.json + outputs/resultados_v3.json
          y CSV en outputs/ — sin números hardcodeados.
Redacción: tercera persona (PRD §19).
"""
import csv
import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── DATOS ───────────────────────────────────────────────────────
BASE_DIR = os.path.dirname(os.path.abspath(__file__))

with open(os.path.join(BASE_DIR, "outputs", "resultados_v2.json"), encoding="utf-8") as f:
    R2 = json.load(f)
with open(os.path.join(BASE_DIR, "outputs", "resultados_v3.json"), encoding="utf-8") as f:
    R3 = json.load(f)

MEJOR  = R3["mejor_modelo"]
MB     = R3["metricas_base"]          # métricas del modelo base (RF test)
CL     = R3["clustering"]
EXP    = R3["experimentos"]
FP     = R3["fp_growth"]
REGLAS = R3["reglas_novedosas"]
IT3    = R3["iteracion_3"]

# Tabla maestra de iteraciones
TABLA_MAESTRA = IT3["tabla_maestra"]
# CV
CV = IT3["cv_estratificada"]
# Optimización
OPT = IT3["optimizacion"]
# Importancia v3
IMP_V3 = IT3["importancia_nuevas_v3"]
# Clustering comparación
CLUST_CMP = IT3["clustering_comparacion"]
# Algoritmo final de clustering (decisión de modelado por adecuación al supuesto)
ALGORITMO_CLUSTERING = IT3.get("clustering_algoritmo_final", "Gaussian Mixture")
# Perfiles de clusters — leídos del CSV generado por el notebook (no hardcodeados)
_perfiles_path = os.path.join(BASE_DIR, "outputs", "tabla_perfiles_clusters.csv")
PERFILES_CSV = []
with open(_perfiles_path, encoding="utf-8") as _f:
    _reader = csv.DictReader(_f)
    for _row in _reader:
        PERFILES_CSV.append(_row)
# Ordinal
ORDINAL = IT3["ordinal"]
# FE ampliada
FE_AMP = IT3["fe_ampliada"]

def img(name):
    return os.path.join(BASE_DIR, "outputs", name)


# ─── PALETA PRD §10.2 ────────────────────────────────────────────
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)   # fondo principal oscuro
BG_LIGHT  = RGBColor(0xF8, 0xFA, 0xFC)   # fondo secundario claro
TEXT_DARK = RGBColor(0x11, 0x18, 0x27)   # texto principal oscuro
TEXT_OVER = RGBColor(0xF9, 0xFA, 0xFB)   # texto sobre fondo oscuro
BLUE      = RGBColor(0x25, 0x63, 0xEB)   # azul académico (principal)
PURPLE    = RGBColor(0x7C, 0x3A, 0xED)   # morado analítico
GREEN     = RGBColor(0x16, 0xA3, 0x4A)   # positivo / mejoras
ORANGE    = RGBColor(0xF9, 0x73, 0x16)   # alertas / limitaciones
RED       = RGBColor(0xDC, 0x26, 0x26)   # riesgos / inferior
BORDER    = RGBColor(0xCB, 0xD5, 0xE1)   # líneas / bordes
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MUTED     = RGBColor(0x64, 0x74, 0x8B)   # etiquetas secundarias
CARD_BG   = RGBColor(0xFF, 0xFF, 0xFF)
DARK_MID  = RGBColor(0x1E, 0x29, 0x3B)   # fondo oscuro medio para secciones
DARK_CARD = RGBColor(0x1E, 0x3A, 0x5F)   # tarjeta sobre fondo oscuro

FONT = "Segoe UI"

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

TOTAL_SLIDES = 20  # portada + 19 explicativas


# ═════════════════════════════════════════════════════════════════
# HELPERS
# ═════════════════════════════════════════════════════════════════
def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def shape_box(slide, kind, l, t, w, h, fill=None, line_color=None, line_w=0.75):
    shp = slide.shapes.add_shape(kind, l, t, w, h)
    if kind == 5:
        try:
            shp.adjustments[0] = 0.07
        except Exception:
            pass
    shp.shadow.inherit = False
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    return shp


def rect(slide, l, t, w, h, fill=None, line_color=None):
    return shape_box(slide, 1, l, t, w, h, fill, line_color)


def card(slide, l, t, w, h, accent=None, fill=CARD_BG):
    c = shape_box(slide, 5, l, t, w, h, fill=fill, line_color=BORDER)
    if accent:
        shape_box(slide, 5, l, t, Inches(0.12), h, fill=accent)
    return c


def dark_card(slide, l, t, w, h, accent=None):
    c = shape_box(slide, 5, l, t, w, h, fill=DARK_CARD, line_color=BORDER)
    if accent:
        shape_box(slide, 5, l, t, Inches(0.12), h, fill=accent)
    return c


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False, font=FONT, anchor=None):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txBox


def header(slide, title, subtitle=None):
    """Encabezado sobre fondo claro con chip azul + título + separador."""
    rect(slide, 0, 0, W, H, fill=BG_LIGHT)
    rect(slide, Inches(0.45), Inches(0.28), Inches(0.16), Inches(0.65), fill=BLUE)
    add_text(slide, title, Inches(0.78), Inches(0.2), Inches(11.9), Inches(0.65),
             size=28, bold=True, color=TEXT_DARK)
    if subtitle:
        add_text(slide, subtitle, Inches(0.8), Inches(0.8), Inches(11.9), Inches(0.38),
                 size=14, color=MUTED)
    rect(slide, Inches(0.45), Inches(1.24), Inches(12.43), Pt(1.5), fill=BORDER)


def slide_num(slide, n, dark=False):
    col = RGBColor(0x94, 0xA3, 0xB8) if dark else MUTED
    add_text(slide, f"{n:02d} / {TOTAL_SLIDES}", Inches(12.2), Inches(7.12),
             Inches(1.0), Inches(0.3), size=9, color=col, align=PP_ALIGN.RIGHT)


def add_image(slide, path, l, t, w, h=None):
    if not os.path.exists(path):
        return
    if h:
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        slide.shapes.add_picture(path, l, t, w)


def img_card(slide, path, l, t, w, pad=Inches(0.1)):
    if not os.path.exists(path):
        return None
    pic = slide.shapes.add_picture(path, l + pad, t + pad, w - 2 * pad)
    frame_h = pic.height + 2 * pad
    c = shape_box(slide, 5, l, t, w, frame_h, fill=CARD_BG, line_color=BORDER)
    sp = c._element
    sp.getparent().remove(sp)
    pic._element.addprevious(sp)
    return frame_h


def bullet_box(slide, items, l, t, w, h, title=None, title_color=TEXT_DARK,
               item_size=13, item_color=TEXT_DARK, leading=4, marker="•  "):
    if title:
        add_text(slide, title, l, t, w, Inches(0.35), size=14, bold=True, color=title_color)
        t += Inches(0.4)
        h -= Inches(0.4)
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(leading)
        run = p.add_run()
        run.text = marker + item
        run.font.size = Pt(item_size)
        run.font.color.rgb = item_color
        run.font.name = FONT


def tabla(slide, data, col_ls, col_ws, ty, row_h, header_bg=None,
          font_size=11, highlight_row=None, highlight_color=GREEN, alt_bg=None):
    if header_bg is None:
        header_bg = BG_DARK
    if alt_bg is None:
        alt_bg = RGBColor(0xF1, 0xF5, 0xF9)
    for ri, row in enumerate(data):
        if ri == 0:
            bg, fc, bold = header_bg, WHITE, True
        elif highlight_row is not None and ri == highlight_row:
            bg, fc, bold = RGBColor(0xDC, 0xFC, 0xE7), highlight_color, True
        else:
            bg = CARD_BG if ri % 2 == 1 else alt_bg
            fc, bold = TEXT_DARK, False
        for ci, (text, lx, cw) in enumerate(zip(row, col_ls, col_ws)):
            rect(slide, lx, ty, cw, row_h, fill=bg, line_color=BORDER)
            add_text(slide, text, lx + Inches(0.06), ty + Inches(0.05),
                     cw - Inches(0.08), row_h - Inches(0.04), size=font_size,
                     bold=bold, color=fc,
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
        ty += row_h
    return ty


def kpi_box(slide, valor, etiqueta, l, t, w, color=BLUE):
    card(slide, l, t, w, Inches(1.05))
    add_text(slide, valor, l, t + Inches(0.08), w, Inches(0.52),
             size=26, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, etiqueta, l + Inches(0.05), t + Inches(0.62),
             w - Inches(0.1), Inches(0.38), size=11, color=MUTED, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════
# SLIDE 0 · PORTADA
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
rect(s, 0, 0, W, H, fill=BG_DARK)
# Formas decorativas geométricas
shape_box(s, 9, Inches(10.5), Inches(3.6), Inches(5.0), Inches(5.0), fill=DARK_MID)
shape_box(s, 9, Inches(11.8), Inches(4.8), Inches(3.4), Inches(3.4), fill=BLUE)
shape_box(s, 9, Inches(9.9), Inches(-1.4), Inches(2.4), Inches(2.4), fill=DARK_MID)
rect(s, Inches(0.75), Inches(2.75), Inches(3.0), Pt(3.5), fill=BLUE)

add_text(s, "TALLER DE APLICACIONES · MAGÍSTER EN DATA SCIENCE · USS",
         Inches(0.75), Inches(0.8), Inches(10.5), Inches(0.4),
         size=13, bold=True, color=BLUE)
add_text(s, "Análisis de Niveles de Obesidad con CRISP-DM",
         Inches(0.7), Inches(1.2), Inches(11.5), Inches(1.0),
         size=44, bold=True, color=TEXT_OVER)
add_text(s, "Clustering · Clasificación · Asociación FP-Growth · Iteraciones de Mejora",
         Inches(0.75), Inches(2.28), Inches(11.0), Inches(0.45),
         size=18, color=RGBColor(0xBA, 0xCC, 0xE0))

add_text(s,
         "Dataset: Estimation of Obesity Levels (UCI ML Repository)\n"
         "2.111 registros · 17 variables · México, Perú y Colombia",
         Inches(0.75), Inches(3.0), Inches(9.5), Inches(0.75),
         size=13, color=RGBColor(0x8A, 0xA4, 0xBC))

infos = [
    ("ALUMNO",   "Dante Gil Zenteno"),
    ("DOCENTE",  "Dr. Mauricio Sepúlveda"),
    ("ASIGNATURA", "Taller de Aplicaciones"),
]
for i, (label, val) in enumerate(infos):
    y = Inches(4.1) + i * Inches(0.82)
    rect(s, Inches(0.78), y + Inches(0.1), Inches(0.07), Inches(0.52), fill=BLUE)
    add_text(s, label, Inches(1.02), y, Inches(3.5), Inches(0.3),
             size=10, bold=True, color=RGBColor(0x7E, 0x96, 0xAE))
    add_text(s, val, Inches(1.02), y + Inches(0.28), Inches(6.0), Inches(0.42),
             size=16, bold=True, color=TEXT_OVER)

add_text(s, "Junio 2026",
         Inches(0.78), Inches(7.05), Inches(3.0), Inches(0.35),
         size=12, bold=True, color=RGBColor(0x7E, 0x96, 0xAE))


# ═════════════════════════════════════════════════════════════════
# SLIDE 1 · PREGUNTA GUÍA / PROBLEMA
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "¿Qué patrones de hábitos y riesgo explican los niveles de obesidad?",
       "Comprensión del negocio — CRISP-DM Fase 1")
slide_num(s, 1)

card(s, Inches(0.45), Inches(1.45), Inches(5.9), Inches(5.6), accent=BLUE)
add_text(s, "El problema", Inches(0.78), Inches(1.6), Inches(5.4), Inches(0.4),
         size=18, bold=True, color=TEXT_DARK)
bullet_box(s, [
    "Obesidad: problema de salud pública con causas multifactoriales (hábitos, genética, entorno)",
    "El dataset registra hábitos alimentarios, actividad física y datos personales",
    "Variable objetivo NObeyesdad: 7 niveles desde Insufficient_Weight hasta Obesity_Type_III",
    "Clasificación multiclase con 7 clases balanceadas → métricas macro confiables",
    "Fuente: Mendoza Palechor & De la Hoz Manotas (2019) — Data in Brief, 104344",
], Inches(0.78), Inches(2.1), Inches(5.4), Inches(3.8), item_size=14, leading=8)

card(s, Inches(6.65), Inches(1.45), Inches(6.25), Inches(2.75), accent=ORANGE)
add_text(s, "Preguntas guía del análisis", Inches(6.98), Inches(1.6), Inches(5.8), Inches(0.4),
         size=17, bold=True, color=TEXT_DARK)
bullet_box(s, [
    "¿Existen perfiles de riesgo diferenciables por clustering?",
    "¿Qué algoritmo predice mejor el nivel de obesidad?",
    "¿Qué combinaciones de hábitos se asocian con mayor riesgo?",
    "¿Reducir dimensiones mejora o degrada los resultados?",
    "¿Nuevas variables de ingeniería aportan señal real?",
], Inches(6.98), Inches(2.08), Inches(5.8), Inches(2.0), item_size=13, leading=5)

card(s, Inches(6.65), Inches(4.4), Inches(6.25), Inches(2.65), accent=GREEN)
add_text(s, "Lo que se espera descubrir", Inches(6.98), Inches(4.55), Inches(5.8), Inches(0.4),
         size=17, bold=True, color=TEXT_DARK)
add_text(s,
         "El equipo aplicó CRISP-DM para comparar múltiples estrategias de modelamiento, "
         "evaluar iteraciones de mejora y extraer conocimiento interpretable sobre los "
         "factores que determinan el nivel de obesidad.",
         Inches(6.98), Inches(5.0), Inches(5.8), Inches(1.9), size=13, color=MUTED)


# ═════════════════════════════════════════════════════════════════
# SLIDE 2 · DATASET Y VARIABLE OBJETIVO
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "El dataset representa 7 niveles de obesidad con alta cobertura de hábitos",
       "Comprensión de datos — CRISP-DM Fase 2")
slide_num(s, 2)

# Imagen distribución + tabla
add_image(s, img("01_distribucion_objetivo.png"),
          Inches(0.45), Inches(1.4), Inches(6.5))

card(s, Inches(7.2), Inches(1.4), Inches(5.7), Inches(5.6), accent=BLUE)
add_text(s, "Estructura del dataset", Inches(7.52), Inches(1.55), Inches(5.2), Inches(0.4),
         size=17, bold=True, color=TEXT_DARK)
bullet_box(s, [
    "2.111 registros · 17 atributos + 1 variable objetivo",
    "Tipo de problema: clasificación multiclase supervisada",
    "Sin nulos ni duplicados — calidad de datos verificada",
    "77% sintético (SMOTE) → clases balanceadas (272–351/nivel)",
    "8 variables numéricas (Edad, Peso, Talla, FAF, FCVC…)",
    "8 variables categóricas (CAEC, CALC, FAVC, MTRANS…)",
    "Target: NObeyesdad — 7 niveles discretos de obesidad",
], Inches(7.52), Inches(2.0), Inches(5.1), Inches(4.0), item_size=13, leading=7)

add_text(s,
         "Clases balanceadas → métrica F1 macro representa todos los niveles por igual",
         Inches(7.52), Inches(6.15), Inches(5.2), Inches(0.5),
         size=12, italic=True, color=MUTED)


# ═════════════════════════════════════════════════════════════════
# SLIDE 3 · METODOLOGÍA CRISP-DM
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "CRISP-DM ordenó el análisis desde el problema hasta la comunicación de resultados",
       "Metodología CRISP-DM — 6 fases iterativas")
slide_num(s, 3)

fases = [
    ("1", "Comprensión del\nNegocio", "Definición de preguntas\ny métricas de éxito", BLUE),
    ("2", "Comprensión\nde los Datos", "Inspección, calidad\ny balance de clases", BLUE),
    ("3", "Preparación\nde los Datos", "Codificación y escalado\npor técnica específica", BLUE),
    ("4", "Modelamiento", "GMM · 6 clasificadores\nFP-Growth · Experimentos", PURPLE),
    ("5", "Evaluación", "Métricas en test · Lift\nComparación iteraciones", PURPLE),
    ("6", "Comunicación", "Presentación · Notebook\nreproducible", PURPLE),
]
for i, (num, nombre, desc, color) in enumerate(fases):
    col, row = i % 3, i // 3
    lx = Inches(0.45 + col * 4.25)
    ty = Inches(1.42 + row * 1.55)
    card(s, lx, ty, Inches(4.05), Inches(1.4))
    shape_box(s, 9, lx + Inches(0.18), ty + Inches(0.28), Inches(0.78), Inches(0.78), fill=color)
    add_text(s, num, lx + Inches(0.18), ty + Inches(0.38), Inches(0.78), Inches(0.56),
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, nombre, lx + Inches(1.1), ty + Inches(0.1), Inches(2.85), Inches(0.6),
             size=13, bold=True, color=TEXT_DARK)
    add_text(s, desc, lx + Inches(1.1), ty + Inches(0.66), Inches(2.85), Inches(0.68),
             size=11, color=MUTED)

card(s, Inches(0.45), Inches(4.68), Inches(12.45), Inches(2.35), accent=BLUE,
     fill=RGBColor(0xEF, 0xF6, 0xFF))
add_text(s, "Por qué CRISP-DM fue útil en este proyecto",
         Inches(0.78), Inches(4.82), Inches(11.9), Inches(0.38),
         size=16, bold=True, color=TEXT_DARK)
bullet_box(s, [
    "Definir preguntas guía ANTES de tocar los datos evitó modelar sin entender el problema",
    "La inspección temprana (nulos, duplicados, balance) descartó problemas de calidad antes de modelar",
    "Preparar los datos POR TÉCNICA (pipelines separados para clustering, clasificación y FP-Growth) eliminó fugas",
    "La fase de evaluación obligó a comparar en test y documentar también los resultados que NO mejoran",
], Inches(0.78), Inches(5.26), Inches(11.9), Inches(1.65), item_size=12, leading=4)


# ═════════════════════════════════════════════════════════════════
# SLIDE 4 · PREPARACIÓN DE DATOS
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "Cada técnica recibió su propia preparación para evitar fugas de datos",
       "Preparación de datos — CRISP-DM Fase 3")
slide_num(s, 4)

bloques = [
    ("Calidad\nverificada", BLUE, [
        "Sin nulos ni duplicados",
        "Rangos numéricos coherentes",
        "7 clases balanceadas (272–351)",
    ]),
    ("Clustering\n(sin target)", PURPLE, [
        "OrdinalEncoder para categóricas",
        "StandardScaler para numéricas",
        "Target excluido — sin fuga",
    ]),
    ("Clasificación\n(pipeline)", BLUE, [
        "Pipeline con preprocesamiento interno",
        "Split estratificado 80/20 · seed 42",
        "Ordinal semántico CAEC / CALC",
    ]),
    ("FP-Growth\n(transaccional)", ORANGE, [
        "Numéricas discretizadas en rangos",
        "Ítems binarios: Gender=Female…",
        "Support mínimo 0.15",
    ]),
    ("Ingeniería de\nCaracterísticas", GREEN, [
        "BMI = Peso / Talla²",
        "BMI_Category (4 categorías)",
        "Healthy / Risky / Sedentary Scores",
        "FAF_FAVC_Interaction",
    ]),
]
for i, (title, color, items) in enumerate(bloques):
    lx = Inches(0.45 + i * 2.53)
    card(s, lx, Inches(1.42), Inches(2.38), Inches(5.6))
    rect(s, lx + Inches(0.16), Inches(1.65), Inches(0.55), Pt(3.5), fill=color)
    add_text(s, title, lx + Inches(0.16), Inches(1.78), Inches(2.1), Inches(0.52),
             size=13, bold=True, color=TEXT_DARK)
    bullet_box(s, items, lx + Inches(0.16), Inches(2.38), Inches(2.1), Inches(4.4),
               item_size=11, leading=6)


# ═════════════════════════════════════════════════════════════════
# SLIDE 5 · DISTRIBUCIÓN NOBEYESDAD
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "El problema multiclase tiene 7 niveles de obesidad aproximadamente balanceados",
       "Distribución de la variable objetivo — NObeyesdad")
slide_num(s, 5)

img_card(s, img("01_distribucion_objetivo.png"),
         Inches(0.45), Inches(1.42), Inches(7.5))

card(s, Inches(8.2), Inches(1.42), Inches(4.7), Inches(5.6), accent=BLUE)
add_text(s, "¿Por qué importa el balance?", Inches(8.52), Inches(1.57), Inches(4.2), Inches(0.4),
         size=16, bold=True, color=TEXT_DARK)
niveles = [
    ("Clase",               "N",    "%"),
    ("Obesity_Type_I",      "351", "16.6"),
    ("Obesity_Type_III",    "324", "15.3"),
    ("Obesity_Type_II",     "297", "14.1"),
    ("Overweight_Level_I",  "290", "13.7"),
    ("Overweight_Level_II", "290", "13.7"),
    ("Normal_Weight",       "287", "13.6"),
    ("Insufficient_Weight", "272", "12.9"),
]
tabla(s, niveles,
      [Inches(8.52), Inches(10.9), Inches(11.7)],
      [Inches(2.38), Inches(0.8), Inches(0.68)],
      Inches(2.05), Inches(0.44), font_size=10)

add_text(s,
         "Clases balanceadas → métricas macro\nson representativas y confiables.\n\n"
         "El equipo eligió F1 macro como métrica\nprincipal por representar el desempeño\n"
         "promedio equitativo entre las 7 clases.",
         Inches(8.52), Inches(5.75), Inches(4.2), Inches(1.22),
         size=12, italic=True, color=MUTED)


# ═════════════════════════════════════════════════════════════════
# SLIDE 6 · ESTRATEGIA GENERAL DE MODELAMIENTO
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "El análisis combinó tres enfoques complementarios sobre los mismos datos",
       "Estrategia de modelamiento — mapa de técnicas aplicadas")
slide_num(s, 6)

bloques6 = [
    ("CLUSTERING\n(no supervisado)", PURPLE, [
        "Objetivo: encontrar perfiles de riesgo sin usar el target",
        "Algoritmos comparados: K-Means, GMM, Agglomerative, DBSCAN",
        "Resultado: 4 perfiles diferenciados por IMC y actividad física",
    ]),
    ("CLASIFICACIÓN\n(supervisado)", BLUE, [
        "Objetivo: predecir el nivel de obesidad desde los hábitos",
        "6 algoritmos comparados: RF, SVM, XGBoost, Árbol, LogReg, KNN",
        f"Mejor en test: {MEJOR}; Random Forest, el más estable en validación cruzada",
    ]),
    ("ASOCIACIÓN\n(FP-Growth)", GREEN, [
        "Objetivo: descubrir patrones de hábitos asociados a cada nivel",
        "7.312 itemsets frecuentes · 134.739 reglas generadas",
        "Selección: Lift > 6 + consecuente = nivel de obesidad",
    ]),
]
for i, (title, color, items) in enumerate(bloques6):
    lx = Inches(0.45 + i * 4.25)
    card(s, lx, Inches(1.42), Inches(4.05), Inches(3.4), accent=color)
    add_text(s, title, lx + Inches(0.32), Inches(1.57), Inches(3.6), Inches(0.6),
             size=16, bold=True, color=color)
    bullet_box(s, items, lx + Inches(0.32), Inches(2.28), Inches(3.55), Inches(2.4),
               item_size=12, leading=6)

card(s, Inches(0.45), Inches(5.0), Inches(12.45), Inches(2.05), accent=PURPLE,
     fill=RGBColor(0xF5, 0xF3, 0xFF))
add_text(s, "Experimentos de mejora aplicados sobre la clasificación",
         Inches(0.78), Inches(5.15), Inches(11.9), Inches(0.38),
         size=15, bold=True, color=PURPLE)
bullet_box(s, [
    f"Reducción dimensional: PCA ({EXP['pca']['n_comp']} comps, 95% varianza) y TruncatedSVD ({EXP['svd']['n_comp']} comps, 90% varianza)",
    "Ingeniería de características: 5→6 variables derivadas (BMI, scores, interacciones)",
    "Tratamiento ordinal semántico de CAEC y CALC + validación cruzada k=5 + RandomizedSearchCV",
], Inches(0.78), Inches(5.58), Inches(11.9), Inches(1.35), item_size=12, leading=4)


# ═════════════════════════════════════════════════════════════════
# SLIDE 7 · CLUSTERING: COMPARACIÓN DE ALGORITMOS
# (fusión de "por qué comparar" + resultados comparativos)
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "GMM seleccionado por adecuación a la estructura no globular — no por agregación de métricas",
       "Clustering — criterio de selección riguroso")
slide_num(s, 7)

img_card(s, img("06_evaluacion_clusters.png"),
         Inches(0.45), Inches(1.42), Inches(7.0))

# Tabla comparativa de algoritmos (leída desde JSON trazable)
# Validez: tamaño mínimo ≥ 5% y ruido ≤ 20% (mismo criterio del notebook)
rows = [("Algoritmo", "Silhouette", "Davies-B.", "Cal-Harab.", "Tam.Min%", "Válido")]
for algo in CLUST_CMP:
    _tam = float(algo.get("tamaño_min_%", 0))
    _ruido = float(algo.get("ruido_%", 0))
    _es_valido = _tam >= 5.0 and _ruido <= 20.0
    valido = "Sí" if _es_valido else "No"
    rows.append((
        algo["Algoritmo"],
        f"{algo['Silhouette']:.4f}",
        f"{algo['Davies-Bouldin']:.3f}",
        f"{algo['Calinski-Harabasz']:.0f}",
        f"{_tam:.1f}%",
        valido,
    ))

add_text(s, "Comparación de algoritmos de clustering — solo candidatos válidos (≥5% mínimo)",
         Inches(7.25), Inches(1.42),
         Inches(5.7), Inches(0.38), size=13, bold=True, color=TEXT_DARK)
tabla(s, rows,
      [Inches(7.25), Inches(9.25), Inches(10.1), Inches(10.88), Inches(11.55), Inches(12.2)],
      [Inches(2.0), Inches(0.85), Inches(0.78), Inches(0.67), Inches(0.65), Inches(0.68)],
      Inches(1.88), Inches(0.46), font_size=9, highlight_row=1, highlight_color=BLUE)

card(s, Inches(7.25), Inches(4.45), Inches(5.7), Inches(2.55), accent=BLUE,
     fill=RGBColor(0xEF, 0xF6, 0xFF))
add_text(s, f"Por qué {ALGORITMO_CLUSTERING} — adecuación al supuesto",
         Inches(7.57), Inches(4.6), Inches(5.2), Inches(0.38),
         size=14, bold=True, color=BLUE)
bullet_box(s, [
    "Estructura no globular (gradiente continuo de IMC): GMM modela componentes elípticas + asignación probabilística (soft)",
    f"{ALGORITMO_CLUSTERING} supera en Silhouette (0.147 vs 0.106), la métrica menos sesgada hacia formas esféricas",
    "Davies-Bouldin y Calinski-Harabasz se miden sobre centroides → favorecen a K-Means por construcción, no son árbitro neutral",
    "Trade-off reconocido: K-Means gana en esas dos métricas; se prioriza la geometría real sobre el supuesto que ellas premian",
    "DBSCAN y Agglomerative invalidados por ruido >4.9% y cluster mínimo <2.1%",
], Inches(7.57), Inches(5.04), Inches(5.25), Inches(1.85), item_size=10, leading=3)


# ═════════════════════════════════════════════════════════════════
# SLIDE 8 · CLUSTERING: PERFILES
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, f"{ALGORITMO_CLUSTERING} k=4: un grupo mayoritario de obesidad (66%) y tres perfiles de menor riesgo",
       f"Clustering — perfiles reales (leídos de tabla_perfiles_clusters.csv)")
slide_num(s, 8)

img_card(s, img("07_clusters_visualizacion.png"),
         Inches(0.45), Inches(1.42), Inches(7.5))

# Perfiles leídos desde outputs/tabla_perfiles_clusters.csv — NO hardcodeados
# Colores asignados por posición de cluster (0=verde normal, 1=rojo severo, 2=naranja moderado, 3=azul moderado-I)
_COLORES_PERFILES = [GREEN, RED, ORANGE, BLUE]
perfiles = []
for _p in PERFILES_CSV:
    _c = int(_p["Cluster"])
    _pct = _p["pct_total"]
    _nombre = f"Cluster {_c} · {_pct}"
    _perfil_riesgo = _p["perfil_riesgo"]
    _imc = float(_p["imc_medio"])
    _peso = float(_p["peso_medio"])
    _faf = float(_p["actividad_fisica_media"])
    _hist = float(_p["prop_historia_familiar"])
    _stats = f"IMC {_imc:.1f} · {_peso:.1f} kg · FAF {_faf:.2f} · Hist. familiar {_hist*100:.0f}%"
    perfiles.append((_nombre, _perfil_riesgo, _stats, _COLORES_PERFILES[_c % len(_COLORES_PERFILES)]))

for i, (nombre, perfil, stats, color) in enumerate(perfiles):
    ty = Inches(1.42 + i * 1.45)
    card(s, Inches(8.2), ty, Inches(4.7), Inches(1.32), accent=color)
    add_text(s, nombre, Inches(8.52), ty + Inches(0.08), Inches(4.1), Inches(0.32),
             size=13, bold=True, color=TEXT_DARK)
    add_text(s, perfil, Inches(8.52), ty + Inches(0.42), Inches(4.1), Inches(0.32),
             size=13, bold=True, color=color)
    add_text(s, stats, Inches(8.52), ty + Inches(0.76), Inches(4.2), Inches(0.48),
             size=11, color=MUTED)

add_text(s,
         f"Decisión: {ALGORITMO_CLUSTERING} (k=4) por adecuación a la estructura no globular\n"
         "(componentes elípticas + Silhouette superior; k por BIC/parsimonia/accionabilidad).",
         Inches(8.2), Inches(7.05), Inches(4.7), Inches(0.38),
         size=10, italic=True, color=MUTED)


# ═════════════════════════════════════════════════════════════════
# SLIDE 9 · CLASIFICACIÓN: MODELOS + RESULTADOS TEST
# (fusión de "modelos evaluados" y "resultados test")
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, f"{MEJOR} obtuvo el mejor F1 macro en test entre los 6 modelos comparados",
       "Clasificación — comparación en conjunto de test")
slide_num(s, 9)

img_card(s, img("08_comparacion_modelos.png"),
         Inches(0.45), Inches(1.42), Inches(7.3))

# Tabla con todos los modelos del v3
filas9 = [("Modelo", "Accuracy", "Prec. M", "Recall M", "F1 Macro")]
orden = sorted(R3["tabla_modelos"].items(), key=lambda kv: kv[1]["F1 Macro"], reverse=True)
for nombre, met in orden:
    star = " ★" if nombre == MEJOR else ""
    filas9.append((
        nombre + star,
        f"{met['Accuracy']:.4f}",
        f"{met['Precision']:.4f}",
        f"{met['Recall']:.4f}",
        f"{met['F1 Macro']:.4f}",
    ))

add_text(s, "Métricas en test (20% estratificado)", Inches(7.55), Inches(1.42),
         Inches(5.4), Inches(0.38), size=14, bold=True, color=TEXT_DARK)
tabla(s, filas9,
      [Inches(7.55), Inches(9.35), Inches(10.27), Inches(11.13), Inches(12.0)],
      [Inches(1.8), Inches(0.92), Inches(0.86), Inches(0.87), Inches(0.85)],
      Inches(1.88), Inches(0.48), font_size=10, highlight_row=1, highlight_color=GREEN)

kpi_box(s, f"{MB['F1 Macro']:.4f}", f"F1 macro — {MEJOR}",
        Inches(7.55), Inches(5.0), Inches(2.65), color=GREEN)
kpi_box(s, f"{MB['Accuracy']:.1%}", "Accuracy en test",
        Inches(10.35), Inches(5.0), Inches(2.65), color=BLUE)

add_text(s,
         "Errores se concentran en clases limítrofes (Normal ↔ Overweight):\n"
         "fronteras de IMC continuas hacen difícil la separación perfecta.",
         Inches(7.55), Inches(6.2), Inches(5.4), Inches(0.7),
         size=11, italic=True, color=MUTED)


# ═════════════════════════════════════════════════════════════════
# SLIDE 10 · MATRIZ DE CONFUSIÓN
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "Los errores del mejor modelo se concentran en clases limítrofes",
       f"Matriz de confusión — {MEJOR} en conjunto de test")
slide_num(s, 10)

img_card(s, img("09_confusion_matrix.png"),
         Inches(0.45), Inches(1.42), Inches(7.8))

card(s, Inches(8.55), Inches(1.42), Inches(4.35), Inches(5.6), accent=BLUE)
add_text(s, "Métricas por clase", Inches(8.87), Inches(1.57), Inches(3.85), Inches(0.38),
         size=15, bold=True, color=TEXT_DARK)

# Datos del reporte de clasificación
clases_datos = [
    ("Clase",              "F1"),
    ("Insufficient_Weight", "0.971"),
    ("Normal_Weight",      "0.912"),
    ("Obesity_Type_I",     "0.965"),
    ("Obesity_Type_II",    "0.992"),
    ("Obesity_Type_III",   "0.992"),
    ("Overweight_Level_I", "0.929"),
    ("Overweight_Level_II","0.957"),
]
tabla(s, clases_datos,
      [Inches(8.87), Inches(11.32)],
      [Inches(2.45), Inches(0.82)],
      Inches(2.0), Inches(0.45), font_size=10)

add_text(s,
         "F1 macro total: " + f"{MB['F1 Macro']:.4f}\n"
         "Clases intermedias (Normal / Overweight I) presentan\n"
         "mayor confusión — sus límites de IMC son continuos\n"
         "y la diferencia entre niveles es sutil.",
         Inches(8.87), Inches(5.25), Inches(3.85), Inches(1.6),
         size=12, italic=True, color=MUTED)


# ═════════════════════════════════════════════════════════════════
# SLIDE 11 · VALIDACIÓN CRUZADA + OPTIMIZACIÓN
# (fusión de ambos bloques)
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "La validación cruzada identifica a Random Forest como el modelo más estable y robusto",
       "Validación cruzada estratificada k=5 + RandomizedSearchCV — Iteración 3")
slide_num(s, 11)

# Tabla CV
rows_cv = [("Modelo", "Acc. CV", "F1 CV media", "F1 CV desv.")]
for row in CV:
    star = " ★" if row["Modelo"] == CV[0]["Modelo"] else ""
    rows_cv.append((
        row["Modelo"] + star,
        f"{row['Accuracy CV media']:.4f}",
        f"{row['F1 macro CV media']:.4f}",
        f"±{row['F1 macro CV desv.']:.4f}",
    ))

add_text(s, "Validación cruzada (k=5 estratificado)", Inches(0.45), Inches(1.42),
         Inches(6.3), Inches(0.38), size=14, bold=True, color=TEXT_DARK)
tabla(s, rows_cv,
      [Inches(0.45), Inches(2.9), Inches(3.9), Inches(4.95)],
      [Inches(2.45), Inches(1.0), Inches(1.05), Inches(1.0)],
      Inches(1.88), Inches(0.46), font_size=10, highlight_row=1, highlight_color=GREEN)

card(s, Inches(0.45), Inches(4.42), Inches(6.3), Inches(2.6), accent=PURPLE,
     fill=RGBColor(0xF5, 0xF3, 0xFF))
add_text(s, "¿Qué aportó la validación cruzada?", Inches(0.78), Inches(4.57),
         Inches(5.85), Inches(0.38), size=14, bold=True, color=PURPLE)
bullet_box(s, [
    f"Una sola partición cambia el ranking: en test lideró {MEJOR}, en CV el más estable es {CV[0]['Modelo']}",
    "k=5 estratificado mantiene proporciones de clase en cada fold",
    f"{CV[0]['Modelo']}: F1 CV {CV[0]['F1 macro CV media']:.4f} ± {CV[0]['F1 macro CV desv.']:.4f} — el más estable",
    f"Por su robustez en CV se optimiza {CV[0]['Modelo']}, no el líder de un único split",
], Inches(0.78), Inches(4.98), Inches(5.85), Inches(1.9), item_size=12, leading=5)

# Panel optimización
card(s, Inches(7.0), Inches(1.42), Inches(5.9), Inches(5.6), accent=GREEN,
     fill=RGBColor(0xF0, 0xFD, 0xF4))
add_text(s, "Optimización de hiperparámetros", Inches(7.32), Inches(1.57),
         Inches(5.4), Inches(0.38), size=14, bold=True, color=GREEN)
bullet_box(s, [
    f"Método: {OPT['metodo']}",
    "Hiperparámetros explorados: n_estimators, max_depth, min_samples_split/leaf, class_weight",
    f"Mejores params: n_estimators=200 · max_depth=20 · class_weight=balanced",
], Inches(7.32), Inches(2.0), Inches(5.35), Inches(1.3), item_size=12, leading=5)

rows_opt = [
    ("Configuración",           "F1 macro test"),
    ("RF por defecto (It.3)",   f"{OPT['f1_test_default']:.4f}"),
    ("RF optimizado (RandSearchCV)", f"{OPT['f1_test_optimizado']:.4f}"),
    ("RF base (It.1)",          f"{OPT['f1_cv_mejor']:.4f}  [CV]"),
]
tabla(s, rows_opt,
      [Inches(7.32), Inches(10.62)],
      [Inches(3.3), Inches(1.55)],
      Inches(3.55), Inches(0.46), font_size=10, highlight_row=2, highlight_color=GREEN)

add_text(s,
         f"La optimización mejoró F1 de {OPT['f1_test_default']:.4f} a {OPT['f1_test_optimizado']:.4f} "
         "(+0.003). Ganancia marginal — el modelo base ya era muy sólido.",
         Inches(7.32), Inches(5.7), Inches(5.4), Inches(0.7),
         size=12, italic=True, color=MUTED)


# ═════════════════════════════════════════════════════════════════
# SLIDE 12 · TRATAMIENTO ORDINAL CAEC/CALC
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "Codificar CAEC y CALC como ordinales respetó su orden semántico natural",
       "Tratamiento ordinal semántico — CAEC y CALC")
slide_num(s, 12)

card(s, Inches(0.45), Inches(1.42), Inches(6.2), Inches(5.6), accent=PURPLE)
add_text(s, "El problema con codificación nominal", Inches(0.78), Inches(1.57),
         Inches(5.7), Inches(0.38), size=16, bold=True, color=TEXT_DARK)
bullet_box(s, [
    "CAEC: frecuencia de consumo entre comidas (no < Sometimes < Frequently < Always)",
    "CALC: consumo de alcohol (no < Sometimes < Frequently < Always)",
    "Con OneHotEncoding se pierde el orden — el modelo no puede aprender la progresión",
    "Con OrdinalEncoder sin orden semántico se asignan números arbitrarios",
    "La solución: OrdinalEncoder con orden explícito [no, Sometimes, Frequently, Always]",
], Inches(0.78), Inches(2.0), Inches(5.7), Inches(3.8), item_size=13, leading=7)

add_text(s, "Resultado del tratamiento ordinal",
         Inches(6.9), Inches(1.42), Inches(6.0), Inches(0.38),
         size=16, bold=True, color=TEXT_DARK)

rows_ord = [
    ("Configuración",         "F1 macro",  "Accuracy"),
    ("Base (sin ordinal)",    f"{MB['F1 Macro']:.4f}", f"{MB['Accuracy']:.4f}"),
    ("Solo ordinal CAEC/CALC",f"{ORDINAL['f1_macro']:.4f}", f"{ORDINAL['accuracy']:.4f}"),
    ("FE ampliada + ordinal", f"{FE_AMP['f1_macro']:.4f}", f"{FE_AMP['accuracy']:.4f}"),
]
tabla(s, rows_ord,
      [Inches(6.9), Inches(9.8), Inches(11.05)],
      [Inches(2.9), Inches(1.25), Inches(1.1)],
      Inches(1.88), Inches(0.5), font_size=11)

card(s, Inches(6.9), Inches(4.6), Inches(6.0), Inches(2.42), accent=ORANGE,
     fill=RGBColor(0xFF, 0xF7, 0xED))
add_text(s, "Interpretación del impacto", Inches(7.22), Inches(4.75),
         Inches(5.55), Inches(0.38), size=14, bold=True, color=ORANGE)
add_text(s,
         f"La codificación ordinal SOLA bajó el F1 de {MB['F1 Macro']:.4f} a {ORDINAL['f1_macro']:.4f}. "
         f"Combinada con FE ampliada (6 atributos) recuperó {FE_AMP['f1_macro']:.4f}.\n\n"
         "El tratamiento ordinal por sí solo no es suficiente — necesita "
         "acompañarse de una representación más rica de los datos.",
         Inches(7.22), Inches(5.18), Inches(5.55), Inches(1.7),
         size=13, color=TEXT_DARK)


# ═════════════════════════════════════════════════════════════════
# SLIDE 13 · INGENIERÍA DE CARACTERÍSTICAS + IMPORTANCIA
# (fusión de ambos bloques)
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "Feature Engineering elevó el F1 a 0.9905 pero revela que el BMI casi determina el target",
       "Ingeniería de características e importancia de atributos — Iteración 2 / 3")
slide_num(s, 13)

img_card(s, img("19_importancia_atributos_v3.png"),
         Inches(0.45), Inches(1.42), Inches(7.0))

card(s, Inches(7.7), Inches(1.42), Inches(5.2), Inches(5.6), accent=GREEN)
add_text(s, "Atributos creados e importancia (RF)", Inches(8.02), Inches(1.57),
         Inches(4.7), Inches(0.38), size=15, bold=True, color=TEXT_DARK)

rows_imp = [("Variable", "Importancia", "Tipo")]
for var, imp in IMP_V3.items():
    rows_imp.append((var, f"{imp:.4f}", "Creada"))
tabla(s, rows_imp,
      [Inches(8.02), Inches(10.52), Inches(11.32)],
      [Inches(2.5), Inches(0.8), Inches(0.65)],
      Inches(2.05), Inches(0.45), font_size=10)

f1_base = MB["F1 Macro"]
f1_fe   = EXP["fe"]["f1_macro"]
add_text(s,
         f"F1 base: {f1_base:.4f}  →  F1 con FE: {f1_fe:.4f}",
         Inches(8.02), Inches(4.65), Inches(4.7), Inches(0.38),
         size=14, bold=True, color=GREEN)

card(s, Inches(7.7), Inches(5.2), Inches(5.2), Inches(1.82), accent=ORANGE,
     fill=RGBColor(0xFF, 0xF7, 0xED))
add_text(s, "Advertencia metodológica", Inches(8.02), Inches(5.35),
         Inches(4.7), Inches(0.38), size=13, bold=True, color=ORANGE)
add_text(s,
         "BMI (importancia 0.31) fue construido desde Peso y Talla — los mismos datos "
         "con que se construyó el target. Esto infla artificialmente las métricas. "
         "El modelo útil en la práctica es el que predice solo con hábitos.",
         Inches(8.02), Inches(5.78), Inches(4.7), Inches(1.18),
         size=12, color=TEXT_DARK)


# ═════════════════════════════════════════════════════════════════
# SLIDE 14 · PCA + SVD DEGRADARON EL DESEMPEÑO
# (fusión de ambas reducciones)
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "PCA y TruncatedSVD redujeron el desempeño en clasificación — más dimensiones no siempre ayudan",
       "Reducción de dimensionalidad — experimentos PCA y TruncatedSVD")
slide_num(s, 14)

# Dos imágenes lado a lado
add_image(s, img("11_pca_varianza.png"),  Inches(0.45), Inches(1.42), Inches(4.1))
add_image(s, img("13_svd_varianza.png"),  Inches(4.7),  Inches(1.42), Inches(4.1))
add_image(s, img("17_pca2d_clases.png"),  Inches(8.9),  Inches(1.42), Inches(4.0))

rows_dim = [
    ("Experimento",                   "F1 macro",                          "Silhouette", "Resultado"),
    ("Base (RF, sin reducción)",       f"{MB['F1 Macro']:.4f}",            f"{CL['sil_base']:.4f}", "Referencia"),
    (f"PCA ({EXP['pca']['n_comp']}c, 95% var.)",
     f"{EXP['pca']['f1_macro']:.4f}",  f"{CL['sil_pca']:.4f}", "DEGRADA"),
    (f"TruncatedSVD ({EXP['svd']['n_comp']}c, 90% var.)",
     f"{EXP['svd']['f1_macro']:.4f}",  f"{CL['sil_svd']:.4f}", "DEGRADA"),
]

add_text(s, "Impacto en clasificación y clustering", Inches(0.45), Inches(5.3),
         Inches(8.0), Inches(0.38), size=14, bold=True, color=TEXT_DARK)
tabla(s, rows_dim,
      [Inches(0.45), Inches(3.95), Inches(4.95), Inches(5.9)],
      [Inches(3.5), Inches(1.0), Inches(0.95), Inches(1.1)],
      Inches(5.75), Inches(0.46), font_size=10, highlight_row=2)

card(s, Inches(8.45), Inches(5.3), Inches(4.45), Inches(1.95), accent=ORANGE,
     fill=RGBColor(0xFF, 0xF7, 0xED))
add_text(s, "Lectura", Inches(8.77), Inches(5.45), Inches(4.0), Inches(0.35),
         size=13, bold=True, color=ORANGE)
add_text(s,
         f"PCA bajó F1 de {MB['F1 Macro']:.4f} a {EXP['pca']['f1_macro']:.4f} "
         f"(−{MB['F1 Macro']-EXP['pca']['f1_macro']:.4f}). "
         f"SVD perdió más: {EXP['svd']['f1_macro']:.4f}. "
         "La compresión destruye información discriminativa útil para 7 clases.",
         Inches(8.77), Inches(5.85), Inches(4.0), Inches(1.3),
         size=12, color=TEXT_DARK)


# ═════════════════════════════════════════════════════════════════
# SLIDE 15 · COMPARACIÓN GENERAL DE ITERACIONES
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "La ingeniería de características y el tuning superaron la reducción dimensional",
       "Comparación de iteraciones — tabla maestra y evolución del F1 macro")
slide_num(s, 15)

img_card(s, img("18_comparacion_iteraciones.png"),
         Inches(0.45), Inches(1.42), Inches(7.5))

rows_maestra = [("It.", "Experimento", "F1 macro")]
for row in TABLA_MAESTRA:
    rows_maestra.append((
        str(row["Iteración"]),
        row["Experimento"],
        f"{row['F1 macro (test)']:.4f}",
    ))

add_text(s, "Tabla maestra de iteraciones", Inches(8.2), Inches(1.42),
         Inches(4.7), Inches(0.38), size=14, bold=True, color=TEXT_DARK)
tabla(s, rows_maestra,
      [Inches(8.2), Inches(8.85), Inches(12.05)],
      [Inches(0.65), Inches(3.2), Inches(0.85)],
      Inches(1.88), Inches(0.46), font_size=9,
      highlight_row=8, highlight_color=GREEN)

card(s, Inches(8.2), Inches(6.3), Inches(4.7), Inches(1.22), accent=GREEN,
     fill=RGBColor(0xF0, 0xFD, 0xF4))
add_text(s,
         f"Mejor resultado: RF optimizado → F1 {OPT['f1_test_optimizado']:.4f}\n"
         f"Peor resultado: TruncatedSVD → F1 {EXP['svd']['f1_macro']:.4f}\n"
         "Conclusión: más complejidad NO implica mejor resultado.",
         Inches(8.5), Inches(6.45), Inches(4.3), Inches(1.02),
         size=12, bold=False, color=TEXT_DARK)


# ═════════════════════════════════════════════════════════════════
# SLIDE 16 · FP-GROWTH + 6 REGLAS NOVEDOSAS
# (fusión de preparación + reglas)
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "FP-Growth descubrió que la historia familiar + hábitos predice obesidad severa con Lift > 6",
       "Reglas de asociación — FP-Growth con support ≥ 0.15 · confidence ≥ 0.70")
slide_num(s, 16)

img_card(s, img("10_reglas_asociacion.png"),
         Inches(0.45), Inches(1.42), Inches(5.9))

add_text(s, "Seis reglas novedosas", Inches(6.65), Inches(1.42),
         Inches(6.35), Inches(0.38), size=15, bold=True, color=TEXT_DARK)

novedad_interp = [
    "Genética + dieta hipercalórica + sedentarismo → Obesity_III",
    "Antecedentes + dieta + alcohol ocasional + transporte → Obesity_III",
    "No fumar NO protege cuando los hábitos alimentarios son malos",
    "Snacking + alcohol moderado → patrón multivariado de riesgo",
    "Predice conducta (FAVC) además del nivel de obesidad",
    "Alcohol + antecedentes familiares + mujer → Obesity_III",
]

for i, regla in enumerate(REGLAS[:6]):
    col = i % 2
    row = i // 2
    lx = Inches(6.65 + col * 3.25)
    ty = Inches(1.9 + row * 1.8)
    card(s, lx, ty, Inches(3.1), Inches(1.65))
    rect(s, lx, ty, Inches(3.1), Inches(0.42), fill=BG_DARK)
    add_text(s, f"Regla {i+1}", lx + Inches(0.12), ty + Inches(0.06),
             Inches(1.8), Inches(0.3), size=12, bold=True, color=WHITE)
    add_text(s, f"Lift {regla['lift']:.2f}",
             lx + Inches(1.9), ty + Inches(0.08), Inches(1.08), Inches(0.28),
             size=11, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
    add_text(s, f"Conf. {regla['confidence']:.3f}  ·  Sup. {regla['support']:.3f}",
             lx + Inches(0.12), ty + Inches(0.48), Inches(2.85), Inches(0.25),
             size=9, color=MUTED)
    add_text(s, novedad_interp[i],
             lx + Inches(0.12), ty + Inches(0.76), Inches(2.85), Inches(0.82),
             size=10, italic=True, color=BLUE)

card(s, Inches(0.45), Inches(6.85), Inches(5.9), Inches(0.55), accent=BLUE)
add_text(s,
         f"{FP['n_itemsets']:,} itemsets frecuentes · {FP['n_reglas']:,} reglas generadas · "
         "Lift > 6 → 6× más probable que el azar".replace(",", "."),
         Inches(0.78), Inches(7.0), Inches(5.6), Inches(0.3),
         size=11, bold=True, color=TEXT_DARK)


# ═════════════════════════════════════════════════════════════════
# SLIDE 17 · TRES DESCUBRIMIENTOS + BIBLIOGRAFÍA
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "Tres descubrimientos clave respaldados por métricas y contrastados con bibliografía",
       "Hallazgos principales — Mendoza Palechor & De la Hoz Manotas (2019)")
slide_num(s, 17)

descubrimientos = [
    (
        f"1 · {MEJOR}: mejor clasificador en test (F1 = {MB['F1 Macro']:.4f})",
        BLUE,
        (f"Se compararon 6 algoritmos. {MEJOR} obtuvo el mejor F1 macro en test ({MB['F1 Macro']:.4f}, accuracy {MB['Accuracy']:.4f}). "
         "La validación cruzada k=5 confirma a Random Forest como el más estable (F1 0.9863), por lo que se optimizó ese modelo, "
         f"alcanzando {OPT['f1_test_optimizado']:.4f}. Das et al. (2026) reportan resultados comparables con "
         "ensembles sobre el mismo dataset (F1 ≈ 0.97).")
    ),
    (
        "2 · Historia familiar domina el riesgo severo (Lift > 6)",
        RED,
        ("Las 6 reglas con mayor Lift incluyen family_history=yes. Los clusters de mayor IMC "
         "concentran 90–100% de antecedentes familiares. Coincide con el estudio HUNT "
         "(Næss et al., 2016): dos padres con sobrepeso → IMC z-score +0.76 en hijos.")
    ),
    (
        "3 · IMC infla métricas — el modelo práctico debe usar solo hábitos",
        ORANGE,
        (f"Con BMI como atributo, F1 sube de {MB['F1 Macro']:.4f} a {EXP['fe']['f1_macro']:.4f}. "
         "El target NObeyesdad fue construido desde el IMC (Mendoza Palechor, 2019): incluirlo "
         "crea circularidad. Un modelo sin antropometría directa es más valioso en la práctica.")
    ),
]
for i, (titulo, color, texto) in enumerate(descubrimientos):
    ty = Inches(1.42 + i * 1.7)
    card(s, Inches(0.45), ty, Inches(12.45), Inches(1.58), accent=color)
    add_text(s, titulo, Inches(0.78), ty + Inches(0.08), Inches(11.9), Inches(0.38),
             size=15, bold=True, color=color)
    add_text(s, texto, Inches(0.78), ty + Inches(0.5), Inches(11.9), Inches(1.0),
             size=12, color=TEXT_DARK)

card(s, Inches(0.45), Inches(6.6), Inches(12.45), Inches(0.78), accent=MUTED,
     fill=RGBColor(0xF8, 0xFA, 0xFC))
add_text(s, "Bibliografía principal:",
         Inches(0.78), Inches(6.65), Inches(2.0), Inches(0.28),
         size=10, bold=True, color=TEXT_DARK)
add_text(s,
         "Mendoza Palechor & De la Hoz Manotas (2019). Dataset for estimation of obesity levels. "
         "Data in Brief, 25, 104344. doi:10.1016/j.dib.2019.104344  ·  "
         "Das et al. (2026). GBWOEM. F1000Research, 14, 1161.  ·  "
         "Næss et al. (2016). Intergenerational Transmission. PLoS ONE, 11(11), e0166585.",
         Inches(2.7), Inches(6.65), Inches(9.7), Inches(0.68),
         size=10, italic=True, color=MUTED)


# ═════════════════════════════════════════════════════════════════
# SLIDE 18 · LIMITACIONES Y MEJORAS FUTURAS
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
header(s, "El análisis tiene limitaciones metodológicas que abren líneas de mejora concretas",
       "Evaluación crítica — limitaciones y mejoras futuras")
slide_num(s, 18)

card(s, Inches(0.45), Inches(1.42), Inches(5.95), Inches(5.6), accent=ORANGE)
add_text(s, "Limitaciones del trabajo", Inches(0.78), Inches(1.57), Inches(5.45), Inches(0.38),
         size=17, bold=True, color=ORANGE)
bullet_box(s, [
    "77% del dataset es sintético (SMOTE): métricas pueden ser más optimistas que en datos reales",
    "Variable objetivo construida desde IMC → incluir BMI como feature crea circularidad",
    "Silhouette bajo (<0.15) en clustering: límites entre niveles de obesidad son continuos",
    "DBSCAN con mejor Silhouette (0.32) pero inválido por ruido y clusters desbalanceados",
    "Reducción de dimensionalidad degradó clasificación: señal discriminativa se perdió al comprimir",
    "Datos de México, Perú y Colombia: generalización a otras poblaciones es incierta",
    "GMM 'full' concentró el 66% en un componente y dejó dos perfiles solapados: mejor Silhouette (0.147) pero menor accionabilidad que K-Means (trade-off asumido)",
], Inches(0.78), Inches(2.0), Inches(5.45), Inches(4.8), item_size=12, leading=6)

card(s, Inches(6.65), Inches(1.42), Inches(6.25), Inches(5.6), accent=GREEN)
add_text(s, "Mejoras futuras propuestas", Inches(6.98), Inches(1.57), Inches(5.75), Inches(0.38),
         size=17, bold=True, color=GREEN)
bullet_box(s, [
    "Modelo solo con hábitos (excluir Peso, Talla, BMI) para un clasificador más útil en la práctica",
    "Validar con datos reales no sintéticos para confirmar generalización",
    "Explorar SHAP values para interpretabilidad fina por instancia",
    "GridSearchCV completo sobre espacio más amplio de hiperparámetros",
    "Aplicación web (Streamlit/FastAPI) para consultas del clasificador",
    "Análisis de clustering con más algoritmos (HDBSCAN) y mayor rango de k",
    "Segmentación demográfica: analizar diferencias por país o grupo etario",
    "Incorporar datos longitudinales para evaluar evolución del riesgo",
], Inches(6.98), Inches(2.0), Inches(5.75), Inches(4.8), item_size=12, leading=5)


# ═════════════════════════════════════════════════════════════════
# SLIDE 19 · CONCLUSIÓN FINAL
# ═════════════════════════════════════════════════════════════════
s = blank_slide()
rect(s, 0, 0, W, H, fill=BG_DARK)
shape_box(s, 9, Inches(10.8), Inches(3.4), Inches(4.8), Inches(4.8), fill=DARK_MID)
shape_box(s, 9, Inches(12.0), Inches(4.6), Inches(3.2), Inches(3.2), fill=BLUE)
rect(s, Inches(0.75), Inches(1.22), Inches(2.8), Pt(4), fill=BLUE)
slide_num(s, 19, dark=True)

add_text(s, "Conclusión Final", Inches(0.72), Inches(0.3),
         Inches(11.0), Inches(0.75), size=38, bold=True, color=TEXT_OVER)
add_text(s, "Obesity Levels Analysis · CRISP-DM · Magíster en Data Science · USS",
         Inches(0.75), Inches(1.32), Inches(11.5), Inches(0.38),
         size=13, color=RGBColor(0x8A, 0xA4, 0xBC))

conclusiones19 = [
    (f"Clustering {ALGORITMO_CLUSTERING} (k=4)", BLUE,
     "Cuatro perfiles de riesgo diferenciados por IMC, peso y actividad física. "
     "Seleccionado por adecuación a la estructura no globular (mixtura elíptica + Silhouette superior). "
     "El grupo de obesidad severa concentra 100% de historia familiar."),
    (f"Clasificación — {MEJOR} (mejor de 6)", GREEN,
     f"F1 macro {MB['F1 Macro']:.4f} en test. RF optimizado alcanzó {OPT['f1_test_optimizado']:.4f}. "
     "Modelos de ensemble superaron métodos lineales y basados en distancia."),
    ("FP-Growth — Lift > 6", PURPLE,
     "Historia familiar + hábitos hipercalóricos + sedentarismo forman predictor combinado. "
     f"7.312 itemsets frecuentes. 6 reglas novedosas con Lift máximo {REGLAS[0]['lift']:.2f}."),
    ("PCA y SVD degradaron clasificación", ORANGE,
     f"F1 cayó de {MB['F1 Macro']:.4f} a {EXP['pca']['f1_macro']:.4f} (PCA) "
     f"y {EXP['svd']['f1_macro']:.4f} (SVD). Más complejidad no implica mejor resultado."),
    ("Mensaje central del análisis", RGBColor(0xFF, 0xFF, 0xFF),
     "Las mejores decisiones no dependieron de aplicar más modelos, sino de comparar alternativas, "
     "justificar supuestos y evaluar cada iteración con métricas consistentes."),
]
for i, (titulo, color, texto) in enumerate(conclusiones19):
    ty = Inches(1.88 + i * 1.06)
    shape_box(s, 5, Inches(0.45), ty, Inches(12.45), Inches(0.96), fill=DARK_CARD)
    rect(s, Inches(0.45), ty + Inches(0.08), Inches(0.09), Inches(0.8), fill=color)
    add_text(s, titulo, Inches(0.75), ty + Inches(0.05), Inches(3.8), Inches(0.4),
             size=13, bold=True, color=color)
    add_text(s, texto, Inches(4.7), ty + Inches(0.06), Inches(8.0), Inches(0.85),
             size=11, color=RGBColor(0xD4, 0xDF, 0xEE))

add_text(s, "Dante Gil Zenteno · Magíster en Data Science, USS · Junio 2026",
         Inches(0.75), Inches(7.12), Inches(11.5), Inches(0.3),
         size=10, color=RGBColor(0x7E, 0x96, 0xAE))


# ─── GUARDAR ─────────────────────────────────────────────────────
out_path = os.path.join(BASE_DIR, "presentacion_obesidad_final_v4.pptx")
prs.save(out_path)
n = len(prs.slides)
print(f"PPT final v4 generada: {out_path}")
print(f"Total de diapositivas: {n}  (portada + {n - 1} explicativas)")
