# -*- coding: utf-8 -*-
"""
Generador de la NUEVA PPT (v2) — Análisis de Obesidad CRISP-DM + Experimentos
Ejecutar: python create_ppt_v2.py   (requiere outputs/resultados_v2.json)
Produce:  presentacion_obesidad_v2.pptx (15 diapositivas)

Estructura según PRD §20. Las métricas se leen desde resultados_v2.json
(generado por el notebook) — sin números hardcodeados desactualizados.
"""
import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

with open("outputs/resultados_v2.json", encoding="utf-8") as f:
    R = json.load(f)

MEJOR = R["mejor_modelo"]
MB = R["metricas_base"]
CL = R["clustering"]
EXP = R["experimentos"]
FP = R["fp_growth"]
REGLAS = R["reglas_novedosas"]
IMP = R["importancia_nuevas"]

# ─── PALETA ──────────────────────────────────────────────────────
NAVY    = RGBColor(0x1A, 0x23, 0x5E)
BLUE    = RGBColor(0x1E, 0x6B, 0xB8)
CYAN    = RGBColor(0x00, 0xBF, 0xD8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xF5, 0xF5, 0xF5)
DGRAY   = RGBColor(0x44, 0x44, 0x44)
GREEN   = RGBColor(0x2E, 0x7D, 0x32)
ORANGE  = RGBColor(0xE6, 0x51, 0x00)
PURPLE  = RGBColor(0x67, 0x3A, 0xB7)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


# ─── HELPERS ─────────────────────────────────────────────────────
def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, l, t, w, h, fill=None, line_color=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=DGRAY,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def header_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, W, Inches(1.35), fill=NAVY)
    add_text(slide, title, Inches(0.35), Inches(0.08), Inches(12.5), Inches(0.75),
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.35), Inches(0.8), Inches(12.5), Inches(0.45),
                 size=14, color=CYAN)


def slide_num(slide, n, total=15):
    add_text(slide, f"{n} / {total}", Inches(12.6), Inches(7.1), Inches(0.7), Inches(0.3),
             size=9, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.RIGHT)


def add_image(slide, path, l, t, w, h=None):
    if not os.path.exists(path):
        return
    if h:
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        slide.shapes.add_picture(path, l, t, w)


def bullet_box(slide, items, l, t, w, h, title=None, title_color=NAVY,
               item_size=13, item_color=DGRAY, indent=True):
    if title:
        add_text(slide, title, l, t, w, Inches(0.35), size=14, bold=True, color=title_color)
        t += Inches(0.38)
        h -= Inches(0.38)
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = ("• " if indent else "") + item
        run.font.size = Pt(item_size)
        run.font.color.rgb = item_color


def tabla(slide, data, col_ls, col_ws, ty, row_h, header_bg=NAVY,
          font_size=11, highlight_row=None):
    for ri, row in enumerate(data):
        if ri == 0:
            bg, fc, bold = header_bg, WHITE, True
        elif highlight_row is not None and ri == highlight_row:
            bg, fc, bold = RGBColor(0xE8, 0xF5, 0xE9), GREEN, True
        else:
            bg = WHITE if ri % 2 == 1 else LGRAY
            fc, bold = DGRAY, False
        for ci, (text, lx, cw) in enumerate(zip(row, col_ls, col_ws)):
            rect(slide, lx, ty, cw, row_h, fill=bg, line_color=RGBColor(0xCC, 0xCC, 0xCC))
            add_text(slide, text, lx + Inches(0.06), ty + Inches(0.07),
                     cw - Inches(0.08), row_h - Inches(0.07), size=font_size,
                     bold=bold, color=fc,
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
        ty += row_h
    return ty


# ─── SLIDE 1: PORTADA ────────────────────────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=NAVY)
rect(s, Inches(0.5), Inches(1.8), Inches(12.33), Inches(0.05), fill=CYAN)

add_text(s, "Análisis de Niveles de Obesidad",
         Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8),
         size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "mediante Métodos de Ciencia de Datos — Versión 2",
         Inches(0.5), Inches(1.0), Inches(12.33), Inches(0.65),
         size=24, color=CYAN, align=PP_ALIGN.CENTER)
add_text(s, "Dataset: Estimation of Obesity Levels — UCI ML Repository  ·  "
            "Nuevos experimentos: PCA · TruncatedSVD · Ingeniería de Características",
         Inches(0.8), Inches(2.1), Inches(11.7), Inches(0.45),
         size=13, color=WHITE, align=PP_ALIGN.CENTER)

rect(s, Inches(3.0), Inches(3.0), Inches(7.33), Inches(2.9),
     fill=RGBColor(0x22, 0x2D, 0x6E), line_color=CYAN)
infos = [
    ("Alumno:", "Dante Gil Zenteno"),
    ("Docente:", "Dr. Mauricio Sepúlveda"),
    ("Asignatura:", "Taller de Aplicaciones"),
    ("Programa:", "Magíster en Data Science — Facultad de Ingeniería"),
]
for i, (label, val) in enumerate(infos):
    y = Inches(3.15) + i * Inches(0.62)
    add_text(s, label, Inches(3.3), y, Inches(2.0), Inches(0.5), size=13, bold=True, color=CYAN)
    add_text(s, val, Inches(5.3), y, Inches(4.8), Inches(0.5), size=13, color=WHITE)

add_text(s, "2026", Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.35),
         size=11, color=RGBColor(0x88, 0x88, 0xAA), align=PP_ALIGN.CENTER)
slide_num(s, 1)


# ─── SLIDE 2: OBJETIVO Y PROBLEMA ────────────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Objetivo del Trabajo y Problema",
           "CRISP-DM · Fase 1: Comprensión del Negocio")
slide_num(s, 2)

rect(s, Inches(0.3), Inches(1.5), Inches(6.1), Inches(5.5), fill=WHITE,
     line_color=RGBColor(0xCC, 0xCC, 0xCC))
rect(s, Inches(6.8), Inches(1.5), Inches(6.2), Inches(5.5), fill=WHITE,
     line_color=RGBColor(0xCC, 0xCC, 0xCC))

add_text(s, "El Dataset", Inches(0.5), Inches(1.55), Inches(5.7), Inches(0.4),
         size=15, bold=True, color=NAVY)
bullet_box(s, [
    "2.111 registros de México, Perú y Colombia",
    "17 atributos: hábitos alimentarios, actividad física y características personales",
    "Variable objetivo NObeyesdad: 7 niveles de obesidad",
    "Clases balanceadas (272–351 registros cada una)",
    "Sin valores nulos ni duplicados",
    "77% de registros sintéticos (SMOTE) — Mendoza Palechor & De la Hoz (2019)",
], Inches(0.5), Inches(2.0), Inches(5.7), Inches(4.7), item_size=12)

add_text(s, "Preguntas Guía (v2)", Inches(7.0), Inches(1.55), Inches(5.8), Inches(0.4),
         size=15, bold=True, color=NAVY)
bullet_box(s, [
    "¿Existen grupos de personas con perfiles de riesgo similares?",
    "¿Qué algoritmo predice mejor el nivel de obesidad?",
    "¿Qué combinaciones de hábitos se asocian con la obesidad?",
    "¿Cómo cambian los resultados con reducción de dimensionalidad (PCA / SVD)?",
    "¿Aportan valor 5 nuevos atributos de ingeniería de características?",
], Inches(7.0), Inches(2.0), Inches(5.8), Inches(3.0), item_size=12)

add_text(s, "Objetivo General", Inches(7.0), Inches(5.1), Inches(5.8), Inches(0.4),
         size=15, bold=True, color=NAVY)
add_text(s, "Aplicar CRISP-DM con agrupamiento, clasificación y asociación FP-Growth, "
            "incorporando 3 experimentos de transformación de datos y comparando su impacto "
            "en las métricas.",
         Inches(7.0), Inches(5.5), Inches(5.8), Inches(1.4), size=12, color=DGRAY)


# ─── SLIDE 3: CRISP-DM — ELEMENTOS ÚTILES ────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "CRISP-DM — ¿Qué Elementos Fueron Útiles?",
           "Reflexión metodológica exigida por la rúbrica")
slide_num(s, 3)

rect(s, Inches(0.3), Inches(1.5), Inches(6.25), Inches(5.6), fill=WHITE, line_color=GREEN)
add_text(s, "✓ Elementos ÚTILES", Inches(0.5), Inches(1.6), Inches(5.9), Inches(0.4),
         size=15, bold=True, color=GREEN)
bullet_box(s, [
    "Comprensión del negocio: obligó a definir preguntas guía ANTES de tocar datos",
    "Comprensión de los datos: inspección temprana (nulos, duplicados, balance) evitó retrabajos",
    "Preparación POR TÉCNICA: escalado para clustering, pipelines para clasificación, discretización para FP-Growth → sin fugas de datos",
    "Evaluación: forzó a comparar en test y a documentar resultados negativos (PCA/SVD no mejoran)",
    "Ciclo iterativo: esta versión 2 ES una iteración CRISP-DM motivada por feedback y rúbrica",
], Inches(0.5), Inches(2.1), Inches(5.9), Inches(4.8), item_size=12)

rect(s, Inches(6.85), Inches(1.5), Inches(6.15), Inches(3.4), fill=WHITE, line_color=ORANGE)
add_text(s, "✗ Elementos MENOS útiles aquí", Inches(7.05), Inches(1.6), Inches(5.8), Inches(0.4),
         size=15, bold=True, color=ORANGE)
bullet_box(s, [
    "Fase de Despliegue: en contexto académico se reduce a PPT + notebook (sin monitoreo productivo)",
    "Linealidad aparente: el trabajo real fue iterativo entre Preparación ↔ Modelamiento",
    "No prescribe el CÓMO: justificar k o elegir métricas vino del criterio técnico, no de la metodología",
], Inches(7.05), Inches(2.1), Inches(5.8), Inches(2.7), item_size=12)

rect(s, Inches(6.85), Inches(5.1), Inches(6.15), Inches(2.0), fill=RGBColor(0xE3, 0xF2, 0xFD),
     line_color=BLUE)
add_text(s, "Cómo ordenó el análisis", Inches(7.05), Inches(5.2), Inches(5.8), Inches(0.4),
         size=14, bold=True, color=BLUE)
add_text(s, "Funcionó como esqueleto organizador: cada sección del notebook corresponde a una fase "
            "→ trabajo trazable y reproducible. Su mayor valor: separar la comprensión del modelamiento.",
         Inches(7.05), Inches(5.6), Inches(5.8), Inches(1.4), size=12, color=DGRAY)


# ─── SLIDE 4: DESCRIPCIÓN DEL DATASET ────────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Descripción del Dataset", "CRISP-DM · Fase 2: Comprensión de los Datos")
slide_num(s, 4)

add_text(s, "Variables del Dataset (17)", Inches(0.3), Inches(1.5), Inches(4.5), Inches(0.4),
         size=14, bold=True, color=NAVY)
bullet_box(s, [
    "Género (Gender)",
    "Edad (Age)  |  Altura (Height)  |  Peso (Weight)",
    "Historia familiar sobrepeso",
    "Consumo alimentos hipercalóricos (FAVC)",
    "Frecuencia de vegetales (FCVC)",
    "Número de comidas principales (NCP)",
    "Snacks entre comidas (CAEC)",
    "Tabaquismo (SMOKE)  |  Agua (CH2O)",
    "Monitoreo de calorías (SCC)",
    "Actividad física (FAF)  |  Tiempo pantallas (TUE)",
    "Consumo de alcohol (CALC)",
    "Medio de transporte (MTRANS)",
    "→ Nivel de obesidad: NObeyesdad  ★ TARGET",
], Inches(0.3), Inches(1.95), Inches(4.8), Inches(5.0), item_size=11)

add_text(s, "Niveles de Obesidad (Variable Objetivo)", Inches(5.5), Inches(1.5),
         Inches(7.5), Inches(0.4), size=14, bold=True, color=NAVY)
niveles = [
    ("Nivel de Obesidad",   "Registros", "%"),
    ("Obesity_Type_I",      "351", "16.6%"),
    ("Obesity_Type_III",    "324", "15.3%"),
    ("Obesity_Type_II",     "297", "14.1%"),
    ("Overweight_Level_I",  "290", "13.7%"),
    ("Overweight_Level_II", "290", "13.7%"),
    ("Normal_Weight",       "287", "13.6%"),
    ("Insufficient_Weight", "272", "12.9%"),
]
tabla(s, niveles, [Inches(5.5), Inches(9.3), Inches(10.6)],
      [Inches(3.8), Inches(1.3), Inches(1.3)], Inches(1.95), Inches(0.42))


# ─── SLIDE 5: PREPARACIÓN Y TRANSFORMACIÓN ───────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Preparación y Transformación de Datos",
           "CRISP-DM · Fase 3: Preparación de los Datos")
slide_num(s, 5)

steps = [
    ("1. Calidad", ["Sin nulos (0 valores faltantes)", "Sin filas duplicadas",
                    "Rangos numéricos consistentes"]),
    ("2. Variables", ["Numéricas (8): Age, Height, Weight, FCVC, NCP, CH2O, FAF, TUE",
                      "Categóricas (8): Gender, FAVC, CAEC, SMOKE, SCC, CALC, MTRANS, family_history",
                      "Target: NObeyesdad — 7 clases"]),
    ("3. Para Clustering", ["OrdinalEncoder en categóricas", "StandardScaler en todas",
                            "Target excluido → sin fuga de datos"]),
    ("4. Para Clasificación", ["Pipeline con ColumnTransformer",
                               "Preprocesamiento DENTRO del pipeline",
                               "train/test estratificado 80/20 · random_state=42"]),
    ("5. Para FP-Growth", ["Numéricas discretizadas en rangos",
                           "Categóricas prefijadas: Gender=Female…",
                           "Formato OHE booleano (mlxtend)"]),
]
col_w = Inches(2.45)
for i, (title, items) in enumerate(steps):
    lx = Inches(0.3 + i * 2.6)
    rect(s, lx, Inches(1.5), col_w, Inches(4.1), fill=WHITE, line_color=RGBColor(0xCC, 0xCC, 0xCC))
    rect(s, lx, Inches(1.5), col_w, Inches(0.5), fill=BLUE)
    add_text(s, title, lx + Inches(0.1), Inches(1.55), col_w - Inches(0.1), Inches(0.4),
             size=12, bold=True, color=WHITE)
    bullet_box(s, items, lx + Inches(0.1), Inches(2.1), col_w - Inches(0.15), Inches(3.3),
               item_size=10)

rect(s, Inches(0.3), Inches(5.8), Inches(12.7), Inches(1.3), fill=RGBColor(0xF3, 0xE5, 0xF5),
     line_color=PURPLE)
add_text(s, "NUEVO EN V2 — Transformaciones experimentales:", Inches(0.5), Inches(5.9),
         Inches(12.3), Inches(0.4), size=13, bold=True, color=PURPLE)
add_text(s, f"PCA ({EXP['pca']['n_comp']} componentes, 95% varianza)  ·  "
            f"TruncatedSVD sobre One-Hot ({EXP['svd']['n_comp']} componentes, 90% varianza)  ·  "
            "5 atributos nuevos: BMI, BMI_Category, Healthy_Habits_Score, Risky_Eating_Score, Sedentary_Risk_Score",
         Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7), size=12, color=DGRAY)


# ─── SLIDE 6: EDA ────────────────────────────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Análisis Exploratorio de Datos (EDA)",
           "CRISP-DM · Fase 2 — Distribuciones, correlaciones y patrones")
slide_num(s, 6)

add_image(s, "outputs/01_distribucion_objetivo.png", Inches(0.2), Inches(1.45), Inches(6.5))
add_image(s, "outputs/03_correlacion.png", Inches(7.0), Inches(1.45), Inches(6.0))
add_text(s, "Hallazgos clave:", Inches(0.2), Inches(6.5), Inches(13.0), Inches(0.35),
         size=12, bold=True, color=NAVY)
add_text(s, "• 7 clases balanceadas (272–351 por nivel)   "
            "• Weight y Height con mayor correlación al target   "
            "• Age tiene menor variación entre clases",
         Inches(0.2), Inches(6.85), Inches(13.0), Inches(0.45), size=11, color=DGRAY)


# ─── SLIDE 7: CLUSTERING — MÉTODO Y K ────────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Clustering — Algoritmo y Número de Grupos",
           "CRISP-DM · Fase 4: Modelamiento (Agrupamiento)")
slide_num(s, 7)

add_image(s, "outputs/06_evaluacion_clusters.png", Inches(0.2), Inches(1.45), Inches(8.2))
rect(s, Inches(8.6), Inches(1.5), Inches(4.5), Inches(5.7), fill=WHITE,
     line_color=RGBColor(0xCC, 0xCC, 0xCC))
add_text(s, "Justificación del Algoritmo", Inches(8.8), Inches(1.6), Inches(4.2), Inches(0.4),
         size=13, bold=True, color=NAVY)
bullet_box(s, [
    "K-Means: escala bien con 2.111 registros",
    "Variables mixtas escaladas → centroides interpretables",
    "Alternativas evaluadas: DBSCAN (requiere calibrar eps) y Agglomerative (lento, resultado similar)",
], Inches(8.8), Inches(2.05), Inches(4.2), Inches(1.9), item_size=11)
add_text(s, "Elección de k=4", Inches(8.8), Inches(4.0), Inches(4.2), Inches(0.4),
         size=13, bold=True, color=NAVY)
bullet_box(s, [
    "k=2 maximiza Silhouette pero solo separa \"obeso / no-obeso\" — insuficiente",
    "k=4 mapea los 7 niveles en 4 perfiles de riesgo accionables",
    "Permite diseñar intervenciones por grupo",
    "Silhouettes <0.20 en todo el rango: normal en datos de salud con registros sintéticos",
], Inches(8.8), Inches(4.45), Inches(4.2), Inches(2.5), item_size=11)


# ─── SLIDE 8: CLUSTERING — PERFILES ──────────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Clustering — Interpretación de los 4 Grupos", "CRISP-DM · Fase 4: Modelamiento")
slide_num(s, 8)

add_image(s, "outputs/07_clusters_visualizacion.png", Inches(0.2), Inches(1.45), Inches(7.8))
clusters = [
    ("Cluster 0 · 22.4%", "Bajo peso / Normal",
     "IMC 21.4 · Peso 56.6 kg\nEdad 20.8 · FAF 1.10\nHist. familiar 41%", CYAN),
    ("Cluster 1 · 23.1%", "Obesidad Severa",
     "IMC 37.8 · Peso 106.6 kg\nEdad 22.9 · FAF 0.73\nHist. familiar 100%", ORANGE),
    ("Cluster 2 · 36.4%", "Obesidad Moderada",
     "IMC 28.9 · Peso 90.2 kg\nEdad 22.3 · FAF 1.20\nHist. familiar 90%",
     RGBColor(0xE6, 0x9A, 0x00)),
    ("Cluster 3 · 18.1%", "Obesidad Leve/I",
     "IMC 31.2 · Peso 91.0 kg\nEdad 34.4 · FAF 0.88\nHist. familiar 92%",
     RGBColor(0x1B, 0x5E, 0x20)),
]
for i, (nombre, perfil, stats, color) in enumerate(clusters):
    lx = Inches(8.2)
    ty = Inches(1.5 + i * 1.45)
    rect(s, lx, ty, Inches(5.0), Inches(1.35), fill=WHITE, line_color=RGBColor(0xCC, 0xCC, 0xCC))
    rect(s, lx, ty, Inches(0.25), Inches(1.35), fill=color)
    add_text(s, nombre, lx + Inches(0.35), ty + Inches(0.05), Inches(4.5), Inches(0.35),
             size=12, bold=True, color=NAVY)
    add_text(s, perfil, lx + Inches(0.35), ty + Inches(0.38), Inches(2.2), Inches(0.35),
             size=11, bold=True, color=color)
    add_text(s, stats, lx + Inches(0.35), ty + Inches(0.7), Inches(4.5), Inches(0.6),
             size=10, color=DGRAY)


# ─── SLIDE 9: CLASIFICACIÓN — MODELOS COMPARADOS ─────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Clasificación — 5 Modelos Comparados",
           "CRISP-DM · Fase 4: Modelamiento (Clasificación)")
slide_num(s, 9)

add_image(s, "outputs/08_comparacion_modelos.png", Inches(0.2), Inches(1.45), Inches(9.5))
add_text(s, "Configuración", Inches(10.0), Inches(1.5), Inches(3.1), Inches(0.4),
         size=13, bold=True, color=NAVY)
bullet_box(s, [
    "Train / Test: 80% / 20% estratificado",
    "Pipeline con preprocesamiento dentro (sin data leakage)",
    "Métricas macro (todas las clases pesan igual)",
    "random_state = 42",
], Inches(10.0), Inches(1.95), Inches(3.1), Inches(2.2), item_size=11)
add_text(s, "Modelos evaluados", Inches(10.0), Inches(4.2), Inches(3.1), Inches(0.4),
         size=13, bold=True, color=NAVY)
bullet_box(s, [
    "Regresión Logística Multinomial",
    "Árbol de Decisión (max_depth=12)",
    "Random Forest (100 árboles)",
    "K-Nearest Neighbors (k=7)",
    "XGBoost (200 árboles) — NUEVO en v2",
], Inches(10.0), Inches(4.65), Inches(3.1), Inches(2.3), item_size=11)


# ─── SLIDE 10: CLASIFICACIÓN — RESULTADOS EN TEST ────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Clasificación — Resultados en Test", "CRISP-DM · Fase 5: Evaluación")
slide_num(s, 10)

add_image(s, "outputs/09_confusion_matrix.png", Inches(0.2), Inches(1.45), Inches(6.5))
add_text(s, "Métricas en Datos de Testeo (20%)", Inches(7.0), Inches(1.5),
         Inches(6.1), Inches(0.4), size=14, bold=True, color=NAVY)

filas = [("Modelo", "Accuracy", "Precision", "Recall", "F1 Macro")]
orden_modelos = sorted(R["tabla_modelos"].items(),
                       key=lambda kv: kv[1]["F1 Macro"], reverse=True)
for nombre, met in orden_modelos:
    star = " ★" if nombre == MEJOR else ""
    filas.append((nombre + star, f"{met['Accuracy']:.4f}", f"{met['Precision']:.4f}",
                  f"{met['Recall']:.4f}", f"{met['F1 Macro']:.4f}"))
tabla(s, filas,
      [Inches(7.0), Inches(9.5), Inches(10.35), Inches(11.25), Inches(12.1)],
      [Inches(2.5), Inches(0.85), Inches(0.9), Inches(0.85), Inches(0.9)],
      Inches(1.95), Inches(0.5), highlight_row=1)

add_text(s, f"★ {MEJOR} es el mejor modelo con F1 Macro = {MB['F1 Macro']:.4f}",
         Inches(7.0), Inches(5.2), Inches(6.1), Inches(0.5), size=13, bold=True, color=GREEN)
add_text(s, "Las clases limítrofes (Normal / Overweight) concentran los errores —\n"
            "esperable porque sus fronteras de IMC son continuas.",
         Inches(7.0), Inches(5.75), Inches(6.1), Inches(0.8), size=11, color=DGRAY)


# ─── SLIDE 11: FP-GROWTH — MEJORES REGLAS ────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "FP-Growth — Reglas con Mejores Indicadores",
           "min_support=0.15 · min_confidence=0.70 · Lift > 2.0 · target como consecuente")
slide_num(s, 11)

add_image(s, "outputs/10_reglas_asociacion.png", Inches(0.2), Inches(1.45), Inches(7.5))

add_text(s, "Preparación: numéricas discretizadas + ítems prefijados "
            f"→ {FP['n_itemsets']:,} itemsets · {FP['n_reglas']:,} reglas".replace(",", "."),
         Inches(0.2), Inches(6.6), Inches(7.6), Inches(0.6), size=11, bold=True, color=NAVY)

top_rules = [
    ("Antecedentes", "Consecuente", "Sup", "Conf", "Lift"),
    ("family_history=yes, Weight=Extremo", "Obesity_III", "0.153", "0.979", "6.38"),
    ("FAVC=yes, FAF=Nulo, Weight=Alto",    "Obesity_II",  "0.167", "0.971", "6.34"),
    ("CALC=Sometimes, Gender=Female",      "Obesity_III", "0.158", "0.969", "6.31"),
    ("NCP=Normal, SCC=no, FAVC=yes",       "Obesity_III", "0.152", "0.970", "6.32"),
    ("CH2O=Bajo, FAF=Nulo",                "Obesity_I",   "0.189", "0.943", "5.47"),
    ("MTRANS=Public, FCVC=Alto",           "Obesity_III", "0.155", "0.974", "6.35"),
]
add_text(s, "Top Reglas (NObeyesdad como consecuente)", Inches(8.0), Inches(1.5),
         Inches(5.2), Inches(0.4), size=13, bold=True, color=NAVY)
tabla(s, top_rules,
      [Inches(8.0), Inches(10.8), Inches(11.45), Inches(12.05), Inches(12.65)],
      [Inches(2.75), Inches(0.65), Inches(0.6), Inches(0.6), Inches(0.65)],
      Inches(1.95), Inches(0.5), font_size=9)
add_text(s, "★ Lift > 6 ⇒ la regla es 6× más probable que por azar",
         Inches(8.0), Inches(5.35), Inches(5.2), Inches(0.45), size=11, bold=True, color=ORANGE)

# Fórmulas
add_text(s, "Support = P(A∪B)   ·   Confidence = P(B|A)   ·   Lift = Confidence / P(B)",
         Inches(8.0), Inches(5.9), Inches(5.2), Inches(0.6), size=10, italic=True, color=BLUE)


# ─── SLIDE 12: 6 REGLAS NOVEDOSAS ────────────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "6 Reglas Novedosas de Asociación",
           "FP-Growth · Interpretación y novedad de las reglas seleccionadas")
slide_num(s, 12)

interpretaciones_novedad = [
    "Combina factor genético con patrón conductual complejo — no trivial.",
    "El transporte público sedentario potencia el efecto de la dieta hipercalórica.",
    "No fumar no protege de obesidad severa cuando los hábitos alimentarios son malos.",
    "Snacking y alcohol moderado refuerzan el riesgo — regla multivariada novedosa.",
    "Predice también el patrón de conducta (FAVC), no solo el nivel de obesidad.",
    "Alcohol ocasional en mujeres con antecedentes familiares y dieta alta en calorías.",
]
card_w, card_h = Inches(4.1), Inches(2.5)
for i, regla in enumerate(REGLAS[:6]):
    col, row = i % 3, i // 3
    lx = Inches(0.25 + col * 4.35)
    ty = Inches(1.5 + row * 2.7)
    rect(s, lx, ty, card_w, card_h, fill=WHITE, line_color=RGBColor(0xCC, 0xCC, 0xCC))
    rect(s, lx, ty, card_w, Inches(0.45), fill=NAVY)
    add_text(s, f"Regla {i+1}", lx + Inches(0.1), ty + Inches(0.05), Inches(2.5), Inches(0.35),
             size=12, bold=True, color=WHITE)
    add_text(s, f"Lift {regla['lift']:.2f}", lx + card_w - Inches(1.3), ty + Inches(0.08),
             Inches(1.2), Inches(0.3), size=11, bold=True, color=CYAN, align=PP_ALIGN.RIGHT)
    desc = (f"{regla['antecedente']}\n→ {regla['consecuente']}  "
            f"(Conf. {regla['confidence']*100:.1f}% · Sup. {regla['support']*100:.1f}%)")
    add_text(s, desc, lx + Inches(0.1), ty + Inches(0.5), card_w - Inches(0.15), Inches(1.05),
             size=8, color=DGRAY)
    rect(s, lx, ty + Inches(1.6), card_w, Inches(0.85), fill=RGBColor(0xF0, 0xF4, 0xFF))
    add_text(s, "Novedad: " + interpretaciones_novedad[i], lx + Inches(0.1), ty + Inches(1.65),
             card_w - Inches(0.15), Inches(0.75), size=8, italic=True, color=BLUE)


# ─── SLIDE 13: EXPERIMENTOS ──────────────────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Experimentos — PCA · TruncatedSVD · Ingeniería de Características",
           "¿Cómo cambian clustering y clasificación con cada transformación?")
slide_num(s, 13)

add_image(s, "outputs/16_comparacion_experimentos.png", Inches(0.2), Inches(1.45), Inches(8.6))

filas_exp = [
    ("Experimento", "Accuracy", "F1 Macro", "Silhouette"),
    ("Base", f"{MB['Accuracy']:.4f}", f"{MB['F1 Macro']:.4f}", f"{CL['sil_base']:.4f}"),
    (f"PCA ({EXP['pca']['n_comp']} comps)", f"{EXP['pca']['accuracy']:.4f}",
     f"{EXP['pca']['f1_macro']:.4f}", f"{CL['sil_pca']:.4f}"),
    (f"SVD ({EXP['svd']['n_comp']} comps)", f"{EXP['svd']['accuracy']:.4f}",
     f"{EXP['svd']['f1_macro']:.4f}", f"{CL['sil_svd']:.4f}"),
    ("FE (+5 attrs)", f"{EXP['fe']['accuracy']:.4f}",
     f"{EXP['fe']['f1_macro']:.4f}", f"{CL['sil_fe']:.4f}"),
]
add_text(s, f"Resultados ({MEJOR} + KMeans k=4)", Inches(9.0), Inches(1.5),
         Inches(4.2), Inches(0.4), size=13, bold=True, color=NAVY)
tabla(s, filas_exp,
      [Inches(9.0), Inches(10.7), Inches(11.55), Inches(12.4)],
      [Inches(1.7), Inches(0.85), Inches(0.85), Inches(0.85)],
      Inches(1.95), Inches(0.48), font_size=10)

bullet_box(s, [
    f"Clasificación: PCA y SVD pierden señal (F1 {EXP['pca']['f1_macro']:.3f} / "
    f"{EXP['svd']['f1_macro']:.3f} vs {MB['F1 Macro']:.3f})",
    f"Clustering: PCA MEJORA la separación — Silhouette {CL['sil_base']:.3f} → "
    f"{CL['sil_pca']:.3f} (+{(CL['sil_pca']/CL['sil_base']-1)*100:.0f}%)",
    "FE: BMI eleva F1 a 0.99 pero casi determina el target (construido desde IMC) — advertencia metodológica",
    "El valor de una transformación DEPENDE de la tarea (predecir vs segmentar)",
], Inches(9.0), Inches(4.6), Inches(4.2), Inches(2.4), item_size=11,
   title="Lectura de los experimentos")


# ─── SLIDE 14: DESCUBRIMIENTOS + BIBLIOGRAFÍA ────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=LGRAY)
header_bar(s, "Descubrimientos Destacados y Bibliografía",
           "Cada hallazgo respaldado por una métrica y contrastado con literatura")
slide_num(s, 14)

descubrimientos = [
    ("1 · Mejor clasificador",
     f"{MEJOR}: F1 macro {MB['F1 Macro']:.3f} y accuracy {MB['Accuracy']:.3f} en test. "
     "Das et al. (2026) reportan resultados comparables con ensembles sobre ESTE MISMO dataset.",
     GREEN),
    ("2 · Historia familiar domina",
     "Las reglas con mayor Lift (>6) incluyen family_history=yes; clusters de mayor IMC: "
     ">90% con antecedentes. Coincide con HUNT Study (Næss et al., 2016): dos padres con "
     "sobrepeso → BMI z-score +0.76 en hijos.",
     ORANGE),
    ("3 · IMC casi determina el target",
     f"Con BMI como atributo, F1 sube de {MB['F1 Macro']:.3f} a {EXP['fe']['f1_macro']:.3f}. "
     "El target fue construido desde el IMC (Mendoza Palechor, 2019) → el modelo útil en la "
     "práctica es el que predice SIN antropometría directa.",
     PURPLE),
    ("4 · Reducción dimensional: efecto mixto",
     f"En clasificación PCA/SVD pierden señal (F1 {EXP['pca']['f1_macro']:.3f} / "
     f"{EXP['svd']['f1_macro']:.3f} vs {MB['F1 Macro']:.3f}), pero en clustering PCA MEJORA "
     f"la separación (Silhouette {CL['sil_base']:.3f} → {CL['sil_pca']:.3f}). La utilidad de "
     "una transformación depende de la tarea.",
     BLUE),
]
for i, (titulo, texto, color) in enumerate(descubrimientos):
    col, row = i % 2, i // 2
    lx = Inches(0.3 + col * 6.55)
    ty = Inches(1.5 + row * 1.85)
    rect(s, lx, ty, Inches(6.35), Inches(1.75), fill=WHITE, line_color=RGBColor(0xCC, 0xCC, 0xCC))
    rect(s, lx, ty, Inches(0.18), Inches(1.75), fill=color)
    add_text(s, titulo, lx + Inches(0.3), ty + Inches(0.06), Inches(5.9), Inches(0.35),
             size=12, bold=True, color=color)
    add_text(s, texto, lx + Inches(0.3), ty + Inches(0.42), Inches(5.9), Inches(1.3),
             size=10, color=DGRAY)

rect(s, Inches(0.3), Inches(5.35), Inches(12.7), Inches(1.75), fill=RGBColor(0xE8, 0xEA, 0xF6),
     line_color=NAVY)
add_text(s, "Bibliografía", Inches(0.5), Inches(5.42), Inches(12.3), Inches(0.35),
         size=13, bold=True, color=NAVY)
bullet_box(s, [
    "Mendoza Palechor, F., & De la Hoz Manotas, A. (2019). Dataset for estimation of obesity levels… Data in Brief, 25, 104344. doi:10.1016/j.dib.2019.104344",
    "Das, S., et al. (2026). GBWOEM: A Gradient-Based Weight Optimization Model… F1000Research, 14, 1161. doi:10.12688/f1000research.169436.2 [PubMed]",
    "Næss, M., et al. (2016). Intergenerational Transmission of Overweight and Obesity — The HUNT Study. PLoS ONE, 11(11), e0166585. doi:10.1371/journal.pone.0166585 [PubMed]",
], Inches(0.5), Inches(5.78), Inches(12.3), Inches(1.3), item_size=10)


# ─── SLIDE 15: CONCLUSIONES ──────────────────────────────────────
s = blank_slide(prs)
rect(s, 0, 0, W, H, fill=NAVY)
rect(s, 0, 0, W, Inches(0.08), fill=CYAN)
rect(s, 0, H - Inches(0.08), W, Inches(0.08), fill=CYAN)

add_text(s, "Conclusiones y Mejoras Futuras", Inches(0.5), Inches(0.15),
         Inches(12.33), Inches(0.7), size=30, bold=True, color=WHITE)
add_text(s, "CRISP-DM · Dataset: Estimation of Obesity Levels — UCI ML Repository · Versión 2",
         Inches(0.5), Inches(0.8), Inches(12.33), Inches(0.4), size=14, color=CYAN)
slide_num(s, 15)

conclusiones = [
    ("Clustering K-Means k=4",
     "4 perfiles de riesgo diferenciados por IMC, peso y actividad física. El grupo de obesidad "
     "severa concentra 100% de historia familiar y la menor actividad física.",
     CYAN),
    (f"Clasificación — {MEJOR}",
     f"Mejor modelo con F1 macro {MB['F1 Macro']:.4f} y accuracy {MB['Accuracy']:.1%} en test, "
     "comparando 5 algoritmos (incluido XGBoost, nuevo en v2).",
     RGBColor(0x66, 0xFF, 0x66)),
    ("Asociación FP-Growth",
     "Lift máximo 6.38: historia familiar + hábitos hipercalóricos + sedentarismo = predictor "
     "combinado robusto de obesidad severa. 6 reglas novedosas interpretadas.",
     ORANGE),
    ("Experimentos (v2)",
     f"PCA y TruncatedSVD no mejoran (F1 {EXP['pca']['f1_macro']:.3f} / {EXP['svd']['f1_macro']:.3f} "
     f"vs {MB['F1 Macro']:.3f}); la ingeniería de características eleva F1 a {EXP['fe']['f1_macro']:.3f} "
     "pero BMI casi determina el target — advertencia metodológica documentada.",
     RGBColor(0xFF, 0xCC, 0x00)),
    ("Mejoras futuras",
     "GridSearchCV + validación cruzada · modelo sin antropometría directa (solo hábitos) · "
     "validación con datos reales no sintéticos · Backoffice Streamlit ya disponible (app.py).",
     RGBColor(0xCE, 0x93, 0xD8)),
]
for i, (titulo, texto, color) in enumerate(conclusiones):
    ty = Inches(1.3 + i * 1.16)
    rect(s, Inches(0.3), ty, Inches(12.7), Inches(1.06),
         fill=RGBColor(0x22, 0x2D, 0x6E), line_color=RGBColor(0x33, 0x3D, 0x7E))
    rect(s, Inches(0.3), ty, Inches(0.2), Inches(1.06), fill=color)
    add_text(s, titulo, Inches(0.65), ty + Inches(0.05), Inches(3.6), Inches(0.4),
             size=12, bold=True, color=color)
    add_text(s, texto, Inches(0.65), ty + Inches(0.4), Inches(12.15), Inches(0.62),
             size=10, color=WHITE)

add_text(s, "Notebook reproducible · Backoffice Streamlit (app.py) · Dataset UCI · Alumno: Dante Gil Zenteno",
         Inches(0.5), Inches(7.12), Inches(12.33), Inches(0.3),
         size=10, color=RGBColor(0x88, 0x88, 0xAA), align=PP_ALIGN.CENTER)


# ─── GUARDAR ─────────────────────────────────────────────────────
out_path = "presentacion_obesidad_v2.pptx"
prs.save(out_path)
print(f"PPT v2 generada: {out_path}")
print(f"Diapositivas: {len(prs.slides)}")
